"""lib/doc_parser.py — Document text extraction for non-PDF formats.

Supported formats:
  - .docx  (Word 2007+, via python-docx)
  - .pptx  (PowerPoint 2007+, via python-pptx — optional)
  - .xlsx  (Excel 2007+, via openpyxl — optional)
  - Plain text (.txt, .md, .csv, .json, .xml, .html, .log, .yaml, .yml, etc.)

All extractors return a dict with:
    text, textLength, totalPages, isScanned, method, warnings
"""

import os

from lib.log import get_logger

logger = get_logger(__name__)

__all__ = ['extract_document_text', 'is_supported_document']

# ── Supported extensions by category ──
_DOCX_EXTS = {'.docx', '.doc'}
_PPTX_EXTS = {'.pptx', '.ppt'}
_XLSX_EXTS = {'.xlsx', '.xls'}
_PLAIN_TEXT_EXTS = {
    '.txt', '.md', '.markdown', '.csv', '.tsv',
    '.json', '.jsonl', '.xml', '.html', '.htm',
    '.log', '.yaml', '.yml', '.toml', '.ini', '.cfg',
    '.rst', '.tex', '.bib', '.srt', '.vtt',
    '.py', '.js', '.ts', '.java', '.c', '.cpp', '.h', '.hpp',
    '.go', '.rs', '.rb', '.php', '.sh', '.bash', '.zsh',
    '.css', '.scss', '.less', '.sql', '.r', '.m', '.swift',
}

_ALL_SUPPORTED = _DOCX_EXTS | _PPTX_EXTS | _XLSX_EXTS | _PLAIN_TEXT_EXTS

# Max chars to extract
_MAX_CHARS = 2_000_000


def is_supported_document(filename: str) -> bool:
    """Check if a filename has a supported document extension."""
    ext = os.path.splitext(filename)[0]  # bug-safe
    ext = os.path.splitext(filename)[1].lower()
    return ext in _ALL_SUPPORTED


def extract_document_text(file_bytes: bytes, filename: str, max_chars: int = 0) -> dict:
    """Extract text from a document file.

    Args:
        file_bytes: Raw file bytes.
        filename: Original filename (used to determine format).
        max_chars: Max chars to extract (0 = unlimited).

    Returns:
        Dict with text, textLength, totalPages, isScanned, method, warnings.
    """
    ext = os.path.splitext(filename)[1].lower()
    limit = max_chars if max_chars > 0 else _MAX_CHARS

    if ext == '.docx':
        return _extract_docx(file_bytes, limit)
    elif ext == '.doc':
        return _extract_doc_legacy(file_bytes, limit)
    elif ext == '.pptx':
        return _extract_pptx(file_bytes, limit)
    elif ext == '.ppt':
        return _extract_ppt_legacy(file_bytes, limit)
    elif ext == '.xlsx':
        return _extract_xlsx(file_bytes, limit)
    elif ext == '.xls':
        return _extract_xls_legacy(file_bytes, limit)
    elif ext in _PLAIN_TEXT_EXTS:
        return _extract_plaintext(file_bytes, filename, limit)
    else:
        return {
            'text': f'[Unsupported format: {ext}]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'unsupported',
            'warnings': [f'Unsupported format: {ext}'],
        }


def _extract_docx(file_bytes: bytes, limit: int) -> dict:
    """Extract text from .docx using python-docx → Markdown-like output."""
    try:
        import docx
    except ImportError:
        logger.warning('[DocParser] python-docx not installed, cannot parse .docx')
        return {
            'text': '[python-docx not installed — run: pip install python-docx]',
            'textLength': 0,
            'totalPages': 1,
            'isScanned': False,
            'method': 'unavailable',
            'warnings': ['python-docx not installed'],
        }

    import io
    warnings = []

    try:
        doc = docx.Document(io.BytesIO(file_bytes))
    except Exception as e:
        logger.error('[DocParser] Failed to open .docx: %s', e, exc_info=True)
        return {
            'text': f'[Failed to parse .docx: {e}]',
            'textLength': 0,
            'totalPages': 1,
            'isScanned': False,
            'method': 'error',
            'warnings': [str(e)],
        }

    parts = []
    total_chars = 0

    # ── Paragraphs with heading detection ──
    for para in doc.paragraphs:
        style_name = (para.style.name or '').lower() if para.style else ''
        text = para.text.strip()
        if not text:
            parts.append('')
            continue

        # Convert Word heading styles to Markdown
        if style_name.startswith('heading'):
            try:
                level = int(style_name.replace('heading', '').strip())
                level = min(max(level, 1), 6)
            except ValueError:
                level = 2
            line = f"{'#' * level} {text}"
        elif style_name in ('title',):
            line = f"# {text}"
        elif style_name in ('subtitle',):
            line = f"## {text}"
        elif style_name.startswith('list'):
            line = f"- {text}"
        else:
            line = text

        total_chars += len(line)
        if total_chars > limit:
            remaining = limit - (total_chars - len(line))
            if remaining > 50:
                parts.append(line[:remaining])
            parts.append(f'\n[…truncated at {limit:,} chars]')
            warnings.append(f'Text truncated at {limit:,} chars')
            break
        parts.append(line)

    # ── Tables → Markdown tables ──
    for table in doc.tables:
        if total_chars > limit:
            break
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace('|', '\\|') for cell in row.cells]
            rows.append('| ' + ' | '.join(cells) + ' |')
        if rows:
            # Insert header separator after first row
            header = rows[0]
            ncols = len(table.rows[0].cells) if table.rows else 1
            separator = '| ' + ' | '.join(['---'] * ncols) + ' |'
            table_md = header + '\n' + separator
            if len(rows) > 1:
                table_md += '\n' + '\n'.join(rows[1:])
            total_chars += len(table_md)
            parts.append('')
            parts.append(table_md)

    text = '\n'.join(parts)
    # Clean up excessive blank lines
    import re
    text = re.sub(r'\n{3,}', '\n\n', text).strip()

    logger.info('[DocParser] Extracted .docx: %d paragraphs, %d tables, %s chars',
                len(doc.paragraphs), len(doc.tables), f'{len(text):,}')

    return {
        'text': text,
        'textLength': len(text),
        'totalPages': max(1, len(text) // 3000),  # rough page estimate
        'isScanned': False,
        'method': 'python-docx',
        'warnings': warnings,
    }


def _extract_pptx(file_bytes: bytes, limit: int) -> dict:
    """Extract text from .pptx using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning('[DocParser] python-pptx not installed, cannot parse .pptx')
        return {
            'text': '[python-pptx not installed — run: pip install python-pptx]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'unavailable',
            'warnings': ['python-pptx not installed'],
        }

    import io
    warnings = []

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        logger.error('[DocParser] Failed to open .pptx: %s', e, exc_info=True)
        return {
            'text': f'[Failed to parse .pptx: {e}]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'error',
            'warnings': [str(e)],
        }

    parts = []
    total_chars = 0
    n_slides = len(prs.slides)

    for si, slide in enumerate(prs.slides, 1):
        slide_parts = [f'## Slide {si}/{n_slides}']
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_parts.append(text)
        slide_text = '\n'.join(slide_parts)
        total_chars += len(slide_text)
        if total_chars > limit:
            parts.append(f'\n[…truncated at slide {si}/{n_slides}]')
            warnings.append(f'Truncated at slide {si}')
            break
        parts.append(slide_text)

    text = '\n\n---\n\n'.join(parts)
    logger.info('[DocParser] Extracted .pptx: %d slides, %s chars',
                n_slides, f'{len(text):,}')

    return {
        'text': text,
        'textLength': len(text),
        'totalPages': n_slides,
        'isScanned': False,
        'method': 'python-pptx',
        'warnings': warnings,
    }


def _extract_xlsx(file_bytes: bytes, limit: int) -> dict:
    """Extract text from .xlsx using openpyxl."""
    try:
        import openpyxl
    except ImportError:
        logger.warning('[DocParser] openpyxl not installed, cannot parse .xlsx')
        return {
            'text': '[openpyxl not installed — run: pip install openpyxl]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'unavailable',
            'warnings': ['openpyxl not installed'],
        }

    import io
    warnings = []

    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    except Exception as e:
        logger.error('[DocParser] Failed to open .xlsx: %s', e, exc_info=True)
        return {
            'text': f'[Failed to parse .xlsx: {e}]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'error',
            'warnings': [str(e)],
        }

    parts = []
    total_chars = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_parts = [f'## Sheet: {sheet_name}']
        rows_data = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            rows_data.append('| ' + ' | '.join(cells) + ' |')
            if len(rows_data) > 1000:
                warnings.append(f'Sheet "{sheet_name}" truncated at 1000 rows')
                break
        if rows_data:
            ncols = max(len(list(ws.iter_rows(values_only=True, max_row=1))[0]) if ws.max_row else 1, 1)
            header = rows_data[0]
            separator = '| ' + ' | '.join(['---'] * ncols) + ' |'
            table_md = header + '\n' + separator
            if len(rows_data) > 1:
                table_md += '\n' + '\n'.join(rows_data[1:])
            sheet_parts.append(table_md)

        sheet_text = '\n'.join(sheet_parts)
        total_chars += len(sheet_text)
        if total_chars > limit:
            parts.append('\n[…truncated]')
            break
        parts.append(sheet_text)

    wb.close()
    text = '\n\n---\n\n'.join(parts)
    logger.info('[DocParser] Extracted .xlsx: %d sheets, %s chars',
                len(wb.sheetnames), f'{len(text):,}')

    return {
        'text': text,
        'textLength': len(text),
        'totalPages': len(wb.sheetnames),
        'isScanned': False,
        'method': 'openpyxl',
        'warnings': warnings,
    }


# ══════════════════════════════════════════════════════
#  Legacy Office formats (.doc, .xls, .ppt)
# ══════════════════════════════════════════════════════

def _extract_doc_legacy(file_bytes: bytes, limit: int) -> dict:
    """Extract text from legacy .doc (Word 97-2003) files.

    Uses olefile to read the raw OLE2 stream and decode Word document text.
    Falls back to basic binary text extraction if olefile is unavailable.
    """
    warnings = []

    # Strategy 1: olefile — read the WordDocument stream
    try:
        import io

        import olefile

        ole = olefile.OleFileIO(io.BytesIO(file_bytes))
        # Word stores text in the 'WordDocument' stream; but the actual plaintext
        # is easier to extract from the '1Table' / '0Table' streams.
        # A simpler approach: read all streams and extract printable text.
        text_parts = []
        for stream_name in ['WordDocument', '1Table', '0Table']:
            if ole.exists(stream_name):
                try:
                    raw = ole.openstream(stream_name).read()
                    # Try UTF-16LE decode (Word's native encoding for text runs)
                    try:
                        decoded = raw.decode('utf-16-le', errors='ignore')
                        # Filter to printable chars
                        cleaned = ''.join(c if c.isprintable() or c in '\n\r\t' else ' ' for c in decoded)
                        if len(cleaned.strip()) > 50:
                            text_parts.append(cleaned)
                    except Exception as e:
                        logger.debug('[DocParser] OLE stream decode failed: %s', e)
                except Exception as e:
                    logger.debug('[DocParser] OLE stream read failed: %s', e)
        ole.close()

        if text_parts:
            # Prefer the longest extracted text
            text = max(text_parts, key=len)
            # Clean up: collapse whitespace runs, normalize line endings
            import re
            text = re.sub(r'[^\S\n]+', ' ', text)
            text = re.sub(r'\n{3,}', '\n\n', text)
            text = text.strip()
            if len(text) > limit:
                text = text[:limit]
                warnings.append(f'Text truncated at {limit:,} chars')
            logger.info('[DocParser] Extracted .doc via olefile: %s chars', f'{len(text):,}')
            return {
                'text': text,
                'textLength': len(text),
                'totalPages': max(1, len(text) // 3000),
                'isScanned': False,
                'method': 'olefile (.doc)',
                'warnings': warnings,
            }
    except ImportError:
        logger.debug('[DocParser] olefile not installed, trying binary fallback for .doc')
    except Exception as e:
        logger.warning('[DocParser] olefile extraction failed for .doc: %s', e)

    # Strategy 2: Binary grep — extract UTF-16LE / ASCII strings from raw bytes
    text = _binary_text_extract(file_bytes, limit)
    if text:
        warnings.append('Extracted via binary text scan (quality may vary)')
        logger.info('[DocParser] Extracted .doc via binary scan: %s chars', f'{len(text):,}')
        return {
            'text': text,
            'textLength': len(text),
            'totalPages': max(1, len(text) // 3000),
            'isScanned': False,
            'method': 'binary-scan (.doc)',
            'warnings': warnings,
        }

    return {
        'text': '[Could not extract text from .doc file — try converting to .docx]',
        'textLength': 0,
        'totalPages': 1,
        'isScanned': False,
        'method': 'unsupported',
        'warnings': ['Legacy .doc text extraction failed'],
    }


def _extract_xls_legacy(file_bytes: bytes, limit: int) -> dict:
    """Extract text from legacy .xls (Excel 97-2003) files using xlrd."""
    try:
        import xlrd
    except ImportError:
        logger.warning('[DocParser] xlrd not installed, cannot parse .xls')
        return {
            'text': '[xlrd not installed — run: pip install xlrd]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'unavailable',
            'warnings': ['xlrd not installed'],
        }

    warnings = []
    try:
        wb = xlrd.open_workbook(file_contents=file_bytes)
    except Exception as e:
        logger.error('[DocParser] Failed to open .xls: %s', e, exc_info=True)
        return {
            'text': f'[Failed to parse .xls: {e}]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'error',
            'warnings': [str(e)],
        }

    parts = []
    total_chars = 0

    for si in range(wb.nsheets):
        ws = wb.sheet_by_index(si)
        sheet_parts = [f'## Sheet: {ws.name}']
        rows_data = []
        for ri in range(min(ws.nrows, 1001)):
            cells = []
            for ci in range(ws.ncols):
                cell = ws.cell(ri, ci)
                if cell.ctype == xlrd.XL_CELL_DATE:
                    try:
                        dt = xlrd.xldate_as_datetime(cell.value, wb.datemode)
                        cells.append(dt.strftime('%Y-%m-%d %H:%M:%S').rstrip(' 00:00:00'))
                    except Exception:
                        cells.append(str(cell.value))
                elif cell.ctype == xlrd.XL_CELL_NUMBER:
                    # Show integers without .0
                    v = cell.value
                    cells.append(str(int(v)) if v == int(v) else str(v))
                elif cell.ctype == xlrd.XL_CELL_BOOLEAN:
                    cells.append('TRUE' if cell.value else 'FALSE')
                else:
                    cells.append(str(cell.value) if cell.value else '')
            rows_data.append('| ' + ' | '.join(c.replace('|', '\\|') for c in cells) + ' |')

        if ws.nrows > 1001:
            warnings.append(f'Sheet "{ws.name}" truncated at 1000 rows')

        if rows_data:
            header = rows_data[0]
            ncols = ws.ncols
            separator = '| ' + ' | '.join(['---'] * ncols) + ' |'
            table_md = header + '\n' + separator
            if len(rows_data) > 1:
                table_md += '\n' + '\n'.join(rows_data[1:])
            sheet_parts.append(table_md)

        sheet_text = '\n'.join(sheet_parts)
        total_chars += len(sheet_text)
        if total_chars > limit:
            parts.append('\n[…truncated]')
            break
        parts.append(sheet_text)

    text = '\n\n---\n\n'.join(parts)
    logger.info('[DocParser] Extracted .xls: %d sheets, %s chars',
                wb.nsheets, f'{len(text):,}')

    return {
        'text': text,
        'textLength': len(text),
        'totalPages': wb.nsheets,
        'isScanned': False,
        'method': 'xlrd (.xls)',
        'warnings': warnings,
    }


def _extract_ppt_legacy(file_bytes: bytes, limit: int) -> dict:
    """Extract text from legacy .ppt (PowerPoint 97-2003) files.

    Uses olefile to read the PowerPoint Document stream and extract text records.
    """
    warnings = []

    try:
        import io
        import struct

        import olefile
    except ImportError:
        logger.warning('[DocParser] olefile not installed, cannot parse .ppt')
        return {
            'text': '[olefile not installed — run: pip install olefile]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'unavailable',
            'warnings': ['olefile not installed'],
        }

    try:
        ole = olefile.OleFileIO(io.BytesIO(file_bytes))
        # PPT stores content in 'PowerPoint Document' stream
        if not ole.exists('PowerPoint Document'):
            ole.close()
            # Fallback to binary extraction
            text = _binary_text_extract(file_bytes, limit)
            if text:
                warnings.append('Extracted via binary text scan')
                return {
                    'text': text,
                    'textLength': len(text),
                    'totalPages': max(1, len(text) // 1500),
                    'isScanned': False,
                    'method': 'binary-scan (.ppt)',
                    'warnings': warnings,
                }
            return {
                'text': '[Could not find PowerPoint content in .ppt file]',
                'textLength': 0,
                'totalPages': 0,
                'isScanned': False,
                'method': 'error',
                'warnings': ['PowerPoint Document stream not found'],
            }

        raw = ole.openstream('PowerPoint Document').read()
        ole.close()

        # Parse PPT binary records — TextBytesAtom (0x0FA8) and TextCharsAtom (0x0FA0)
        # contain the actual slide text.
        text_parts = []
        offset = 0
        while offset < len(raw) - 8:
            struct.unpack_from('<H', raw, offset)[0]
            rec_type = struct.unpack_from('<H', raw, offset + 2)[0]
            rec_len = struct.unpack_from('<I', raw, offset + 4)[0]
            offset += 8
            if offset + rec_len > len(raw):
                break
            if rec_type == 0x0FA0:  # TextCharsAtom — UTF-16LE text
                try:
                    text_parts.append(raw[offset:offset + rec_len].decode('utf-16-le', errors='ignore'))
                except Exception as e:
                    logger.debug('[DocParser] PPT TextCharsAtom decode failed: %s', e)
            elif rec_type == 0x0FA8:  # TextBytesAtom — ASCII/Latin-1 text
                try:
                    text_parts.append(raw[offset:offset + rec_len].decode('latin-1', errors='ignore'))
                except Exception as e:
                    logger.debug('[DocParser] PPT TextBytesAtom decode failed: %s', e)
            offset += rec_len

        if text_parts:
            # Join with newlines, clean up
            import re
            text = '\n'.join(t.strip() for t in text_parts if t.strip())
            text = re.sub(r'\n{3,}', '\n\n', text).strip()
            if len(text) > limit:
                text = text[:limit]
                warnings.append(f'Text truncated at {limit:,} chars')
            n_slides = max(1, text.count('\n\n') + 1)  # rough estimate
            logger.info('[DocParser] Extracted .ppt: ~%d text blocks, %s chars',
                        len(text_parts), f'{len(text):,}')
            return {
                'text': text,
                'textLength': len(text),
                'totalPages': n_slides,
                'isScanned': False,
                'method': 'olefile (.ppt)',
                'warnings': warnings,
            }

    except Exception as e:
        logger.warning('[DocParser] olefile extraction failed for .ppt: %s', e)

    # Fallback to binary extraction
    text = _binary_text_extract(file_bytes, limit)
    if text:
        warnings.append('Extracted via binary text scan (quality may vary)')
        return {
            'text': text,
            'textLength': len(text),
            'totalPages': max(1, len(text) // 1500),
            'isScanned': False,
            'method': 'binary-scan (.ppt)',
            'warnings': warnings,
        }

    return {
        'text': '[Could not extract text from .ppt file — try converting to .pptx]',
        'textLength': 0,
        'totalPages': 0,
        'isScanned': False,
        'method': 'unsupported',
        'warnings': ['Legacy .ppt text extraction failed'],
    }


def _binary_text_extract(file_bytes: bytes, limit: int) -> str:
    """Last-resort text extraction from binary Office files.

    Scans the raw bytes for UTF-16LE and ASCII text runs,
    filters to printable content, and returns the best result.
    """
    import re

    # Extract UTF-16LE strings (≥6 chars / 12 bytes)
    utf16_pattern = re.compile(rb'(?:[\x20-\x7e]\x00){6,}')
    matches = utf16_pattern.findall(file_bytes[:limit * 3])
    utf16_text = ''.join(
        m.decode('utf-16-le', errors='ignore') for m in matches
    )

    # Extract ASCII strings (≥8 chars)
    ascii_pattern = re.compile(rb'[\x20-\x7e]{8,}')
    matches = ascii_pattern.findall(file_bytes[:limit * 2])
    ascii_text = ''.join(
        m.decode('ascii', errors='ignore') + '\n' for m in matches
    )

    # Pick the longer, more useful result
    text = utf16_text if len(utf16_text) > len(ascii_text) else ascii_text
    text = re.sub(r'[^\S\n]+', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text).strip()

    if len(text) > limit:
        text = text[:limit]

    return text if len(text) > 50 else ''


def _extract_plaintext(file_bytes: bytes, filename: str, limit: int) -> dict:
    """Extract text from plain-text files with encoding detection."""
    warnings = []
    text = None

    # Try UTF-8 first, then common fallbacks
    for encoding in ('utf-8', 'utf-8-sig', 'gbk', 'gb18030', 'latin-1'):
        try:
            text = file_bytes.decode(encoding)
            if encoding not in ('utf-8', 'utf-8-sig'):
                logger.debug('[DocParser] Decoded %s with %s', filename, encoding)
            break
        except (UnicodeDecodeError, LookupError):
            continue

    if text is None:
        # Last resort: lossy decode
        text = file_bytes.decode('utf-8', errors='replace')
        warnings.append('File contains non-UTF-8 characters (lossy decode)')

    if len(text) > limit:
        text = text[:limit]
        warnings.append(f'Text truncated at {limit:,} chars')

    ext = os.path.splitext(filename)[1].lower()
    logger.info('[DocParser] Extracted plaintext %s (%s): %s chars',
                filename, ext, f'{len(text):,}')

    return {
        'text': text,
        'textLength': len(text),
        'totalPages': 1,
        'isScanned': False,
        'method': f'plaintext ({ext})',
        'warnings': warnings,
    }
