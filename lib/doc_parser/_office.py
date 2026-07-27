"""lib/doc_parser/_office.py — Modern OOXML (Office 2007+) text extractors.

Provides:
  - _extract_docx  (Word 2007+, via python-docx)
  - _extract_pptx  (PowerPoint 2007+, via python-pptx)
  - _extract_xlsx  (Excel 2007+, via openpyxl)

All optional dependencies are imported lazily inside each extractor so that
importing this module never hard-fails when a backend package is missing.
"""

from lib.log import get_logger

from lib.doc_parser._truncation import truncation_warning

logger = get_logger(__name__)

# ── .xlsx scan bounds ──
# Guard against grossly-inflated worksheet dimensions (common with embedded
# images / drawing anchors): cap the rows and columns we iterate, and bail out
# of long runs of fully-empty rows instead of walking to a phantom max_row.
_XLSX_MAX_ROWS = 1000
_XLSX_MAX_COLS = 200
_XLSX_MAX_EMPTY_RUN = 50


def _extract_docx(file_bytes: bytes, limit: int) -> dict:
    """Extract text from .docx using python-docx → Markdown-like output."""
    try:
        import docx
    except ImportError:
        logger.warning('[DocParser] python-docx not installed, cannot parse .docx')
        return {
            'text': '[python-docx not installed — run: pip install python-docx]',
            'textLength': 0,
            'totalPages': 1,
            'isScanned': False,
            'method': 'unavailable',
            'warnings': ['python-docx not installed'],
        }

    import io
    warnings = []

    try:
        doc = docx.Document(io.BytesIO(file_bytes))
    except Exception as e:
        logger.error('[DocParser] Failed to open .docx: %s', e, exc_info=True)
        return {
            'text': f'[Failed to parse .docx: {e}]',
            'textLength': 0,
            'totalPages': 1,
            'isScanned': False,
            'method': 'error',
            'warnings': [str(e)],
        }

    parts = []
    total_chars = 0
    # Denominator for any truncation warning: the full text length this doc
    # WOULD have produced. Computed from the paragraph texts already in
    # memory via python-docx, so a cut can say "kept N of M chars" rather
    # than only naming the limit it hit.
    doc_total_chars = sum(len(p.text.strip()) for p in doc.paragraphs)

    # ── Paragraphs with heading detection ──
    for para in doc.paragraphs:
        style_name = (para.style.name or '').lower() if para.style else ''
        text = para.text.strip()
        if not text:
            parts.append('')
            continue

        # Convert Word heading styles to Markdown
        if style_name.startswith('heading'):
            try:
                level = int(style_name.replace('heading', '').strip())
                level = min(max(level, 1), 6)
            except ValueError as _e_audit:
                logger.debug('[doc_parser] _extract_docx caught %s: %s', type(_e_audit).__name__, _e_audit)
                level = 2
            line = f"{'#' * level} {text}"
        elif style_name in ('title',):
            line = f"# {text}"
        elif style_name in ('subtitle',):
            line = f"## {text}"
        elif style_name.startswith('list'):
            line = f"- {text}"
        else:
            line = text

        total_chars += len(line)
        if total_chars > limit:
            remaining = limit - (total_chars - len(line))
            if remaining > 50:
                parts.append(line[:remaining])
            parts.append(f'\n[…truncated at {limit:,} chars]')
            warnings.append(truncation_warning(
                kept=total_chars - len(line) + max(remaining, 0),
                total=doc_total_chars, unit='chars',
                detail=f'char limit {limit:,}'))
            break
        parts.append(line)

    # ── Tables → Markdown tables ──
    for table in doc.tables:
        if total_chars > limit:
            break
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace('|', '\\|') for cell in row.cells]
            rows.append('| ' + ' | '.join(cells) + ' |')
        if rows:
            # Insert header separator after first row
            header = rows[0]
            ncols = len(table.rows[0].cells) if table.rows else 1
            separator = '| ' + ' | '.join(['---'] * ncols) + ' |'
            table_md = header + '\n' + separator
            if len(rows) > 1:
                table_md += '\n' + '\n'.join(rows[1:])
            total_chars += len(table_md)
            parts.append('')
            parts.append(table_md)

    text = '\n'.join(parts)
    # Clean up excessive blank lines
    import re
    text = re.sub(r'\n{3,}', '\n\n', text).strip()

    logger.info('[DocParser] Extracted .docx: %d paragraphs, %d tables, %s chars',
                len(doc.paragraphs), len(doc.tables), f'{len(text):,}')

    return {
        'text': text,
        'textLength': len(text),
        'totalPages': max(1, len(text) // 3000),  # rough page estimate
        'isScanned': False,
        'method': 'python-docx',
        'warnings': warnings,
    }


def _extract_pptx(file_bytes: bytes, limit: int) -> dict:
    """Extract text from .pptx using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning('[DocParser] python-pptx not installed, cannot parse .pptx')
        return {
            'text': '[python-pptx not installed — run: pip install python-pptx]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'unavailable',
            'warnings': ['python-pptx not installed'],
        }

    import io
    warnings = []

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        logger.error('[DocParser] Failed to open .pptx: %s', e, exc_info=True)
        return {
            'text': f'[Failed to parse .pptx: {e}]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'error',
            'warnings': [str(e)],
        }

    parts = []
    total_chars = 0
    n_slides = len(prs.slides)

    for si, slide in enumerate(prs.slides, 1):
        slide_parts = [f'## Slide {si}/{n_slides}']
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_parts.append(text)
        slide_text = '\n'.join(slide_parts)
        total_chars += len(slide_text)
        if total_chars > limit:
            parts.append(f'\n[…truncated at slide {si}/{n_slides}]')
            warnings.append(truncation_warning(
                kept=si - 1, total=n_slides, unit='slides',
                detail=f'stopped at slide {si}'))
            break
        parts.append(slide_text)

    text = '\n\n---\n\n'.join(parts)
    logger.info('[DocParser] Extracted .pptx: %d slides, %s chars',
                n_slides, f'{len(text):,}')

    return {
        'text': text,
        'textLength': len(text),
        'totalPages': n_slides,
        'isScanned': False,
        'method': 'python-pptx',
        'warnings': warnings,
    }


def _extract_xlsx(file_bytes: bytes, limit: int) -> dict:
    """Extract text from .xlsx using openpyxl."""
    try:
        import openpyxl
    except ImportError:
        logger.warning('[DocParser] openpyxl not installed, cannot parse .xlsx')
        return {
            'text': '[openpyxl not installed — run: pip install openpyxl]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'unavailable',
            'warnings': ['openpyxl not installed'],
        }

    import io
    warnings = []

    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    except Exception as e:
        logger.error('[DocParser] Failed to open .xlsx: %s', e, exc_info=True)
        return {
            'text': f'[Failed to parse .xlsx: {e}]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'error',
            'warnings': [str(e)],
        }

    parts = []
    total_chars = 0
    n_sheets = len(wb.sheetnames)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_parts = [f'## Sheet: {sheet_name}']

        # Worksheet dimensions are often grossly inflated — embedded images,
        # drawing anchors, or stray formatting can push max_row/max_column to
        # tens of thousands even when the real data is a handful of cells.
        # Iterating the full reported range would emit millions of empty
        # cells (slow + useless output), so bound the scan up front and trim
        # trailing-empty cells / skip empty rows as we go.
        col_cap = min(ws.max_column or _XLSX_MAX_COLS, _XLSX_MAX_COLS)

        rows_data = []
        n_real_cols = 0
        empty_run = 0
        truncated_rows = False
        # ★ Every cut must be able to report its DENOMINATOR. A warning that
        # says "truncated at 1000 rows" without saying "of 5000" gives the
        # model a numerator with no scale — it cannot tell 20% from 99%.
        rows_scanned = 0          # data rows actually walked (excl. blanks)
        empty_run_stopped_at = 0  # row index where a blank run ended the scan
        for row in ws.iter_rows(values_only=True, max_col=col_cap):
            cells = list(row)
            while cells and cells[-1] is None:
                cells.pop()
            if not cells:
                empty_run += 1
                if empty_run > _XLSX_MAX_EMPTY_RUN:
                    # This break used to be entirely SILENT. A sheet shaped
                    # "summary block / 60 blank rows / detail block" lost the
                    # whole detail block with no trace in the output at all —
                    # worse than the row cap, which at least admits it fired.
                    empty_run_stopped_at = rows_scanned
                    break
                continue
            empty_run = 0
            rows_scanned += 1
            n_real_cols = max(n_real_cols, len(cells))
            rows_data.append(
                '| ' + ' | '.join(
                    (str(c).replace('|', '\\|') if c is not None else '') for c in cells
                ) + ' |'
            )
            if len(rows_data) >= _XLSX_MAX_ROWS:
                truncated_rows = True
                break

        # Sheet dimensions as reported by the workbook — the denominator the
        # caller needs. Guarded because max_row/max_column can be None.
        sheet_rows = ws.max_row or 0
        sheet_cols = ws.max_column or 0

        if truncated_rows:
            warnings.append(truncation_warning(
                kept=len(rows_data), total=sheet_rows, unit='rows',
                scope=f'Sheet "{sheet_name}"',
                detail=f'row cap {_XLSX_MAX_ROWS:,}'))
        if empty_run_stopped_at:
            warnings.append(truncation_warning(
                kept=empty_run_stopped_at, total=sheet_rows, unit='rows',
                scope=f'Sheet "{sheet_name}"',
                detail=(f'stopped after {_XLSX_MAX_EMPTY_RUN} consecutive '
                        f'blank rows — content below a long blank gap is '
                        f'missing')))
        if sheet_cols > _XLSX_MAX_COLS:
            warnings.append(truncation_warning(
                kept=_XLSX_MAX_COLS, total=sheet_cols, unit='columns',
                scope=f'Sheet "{sheet_name}"'))

        if rows_data:
            ncols = max(n_real_cols, 1)
            header = rows_data[0]
            separator = '| ' + ' | '.join(['---'] * ncols) + ' |'
            table_md = header + '\n' + separator
            if len(rows_data) > 1:
                table_md += '\n' + '\n'.join(rows_data[1:])
            sheet_parts.append(table_md)

        sheet_text = '\n'.join(sheet_parts)
        total_chars += len(sheet_text)
        if total_chars > limit:
            parts.append('\n[…truncated]')
            break
        parts.append(sheet_text)

    wb.close()
    text = '\n\n---\n\n'.join(parts)
    logger.info('[DocParser] Extracted .xlsx: %d sheets, %s chars',
                n_sheets, f'{len(text):,}')

    return {
        'text': text,
        'textLength': len(text),
        'totalPages': len(wb.sheetnames),
        'isScanned': False,
        'method': 'openpyxl',
        'warnings': warnings,
    }
