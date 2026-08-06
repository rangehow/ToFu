"""lib/slides/export_pptx.py — PPTD → PPTX native writer (python-pptx).

The deliverable path (docs/SLIDES_CAPABILITY_DESIGN.md §4.4): every element
becomes REAL OOXML — text boxes keep their runs, shapes stay shapes, tables
stay tables — so the deck opens editable in PowerPoint / WPS. What v1
deliberately rasterises: icons (SVG glyph → transparent PNG via headless
Chrome) and anything chart-like (the DSL v1 builds those from shapes anyway).

Mapping contract (1 px = 1 pt = 12700 EMU):

  * text   → add_textbox; paragraphs/runs from lib.slides.richtext;
             letterSpacing → rPr@spc; lineHeight multiple → line_spacing;
             CJK family set on BOTH latin and ea typeface slots;
  * shape  → MSO_SHAPE autoshape + adjustments; gradient/alpha fills and
             shadows are injected XML (python-pptx has no API for them);
  * line   → freeform polyline (smooth curves are Catmull-Rom-sampled);
             arrowheads via a:headEnd/a:tailEnd;
  * image  → add_picture; crop+fit resolved to crop fractions (cover) or a
             computed placement rect (contain); cropShape → prstGeom;
  * table  → add_table; merged cells via cell.merge; borders via tcPr XML;
  * icon   → rasterised PNG (cached in the design asset store).

Every slide gets a fade transition post-written into the slide XML — the
CT_Slide child order (cSld → clrMapOvr → transition → timing/extLst) is
validated after the rewrite, because Office silently ignores a misplaced
transition.
"""

from __future__ import annotations

import os
import re
import zipfile

from lib.log import get_logger
from lib.slides.pptd import (Deck, Page, resolve_media, table_style,
                             text_style)
from lib.slides.richtext import parse_rich_text

logger = get_logger(__name__)

__all__ = ['export_pptx', 'ExportError']

EMU_PER_PT = 12700


class ExportError(RuntimeError):
    pass


def _emu(v: float) -> int:
    return int(round(float(v) * EMU_PER_PT))


def _rgb(hex_color: str) -> tuple:
    """'#RRGGBB[AA]' → (RGBColor, alpha 0..1)."""
    from pptx.dml.color import RGBColor
    v = (hex_color or '#000000').strip()
    if not v.startswith('#'):
        v = '#000000'
    alpha = 1.0
    if len(v) == 9:
        alpha = int(v[7:9], 16) / 255.0
        v = v[:7]
    if not re.match(r'^#[0-9a-fA-F]{6}$', v):
        v = '#000000'
    return RGBColor.from_string(v[1:]), alpha


# ── XML helpers (the bits python-pptx has no API for) ─────

def _set_alpha_on_fill(fill, alpha: float) -> None:
    if alpha >= 0.999:
        return
    from pptx.oxml.ns import qn
    for tag in ('a:solidFill',):
        solid = fill._xPr.find(qn(tag))
        if solid is None:
            continue
        srgb = solid.find(qn('a:srgbClr'))
        if srgb is not None:
            el = srgb.makeelement(qn('a:alpha'),
                                  {'val': str(int(alpha * 100000))})
            srgb.append(el)


def _gradient_fill_xml(spPr, fill: dict, theme: dict) -> None:
    """Replace spPr's fill with a gradFill built from the PPTD gradient."""
    from pptx.oxml.ns import qn
    from lxml import etree

    _NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    grad = etree.SubElement(spPr, f'{{{_NS}}}gradFill')
    gs_lst = etree.SubElement(grad, f'{{{_NS}}}gsLst')
    for stop in fill.get('stops') or []:
        if not isinstance(stop, dict):
            continue
        gs = etree.SubElement(gs_lst, f'{{{_NS}}}gs')
        gs.set('pos', str(int(max(0.0, min(1.0,
                               float(stop.get('position', 0)))) * 100000)))
        color = stop.get('color') or '#000000'
        from lib.slides.pptd import resolve_color
        color = resolve_color(color, theme, '#000000')
        alpha = 1.0
        if re.match(r'^#[0-9a-fA-F]{8}$', color):
            alpha = int(color[7:9], 16) / 255.0
            color = color[:7]
        srgb = etree.SubElement(gs, f'{{{_NS}}}srgbClr')
        srgb.set('val', color.lstrip('#'))
        if alpha < 0.999:
            a = etree.SubElement(srgb, f'{{{_NS}}}alpha')
            a.set('val', str(int(alpha * 100000)))
    if fill.get('gradientType') == 'radial':
        path = etree.SubElement(grad, f'{{{_NS}}}path')
        path.set('path', 'circle')
        ftr = etree.SubElement(path, f'{{{_NS}}}fillToRect')
        for k, v in (('l', '50000'), ('t', '50000'), ('r', '50000'),
                     ('b', '50000')):
            ftr.set(k, v)
    else:
        # OOXML lin ang: 60000ths of a degree, 0 = left→right, clockwise —
        # the same convention PPTD uses, so the angle passes straight through.
        lin = etree.SubElement(grad, f'{{{_NS}}}lin')
        lin.set('ang', str(int(float(fill.get('angle', 0)) * 60000)))
        lin.set('scaled', '1')
    # gradFill must precede a:ln — move it before any line element.
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        spPr.remove(grad)
        ln.addprevious(grad)


def _set_run_font(run, family: str) -> None:
    """Set latin AND east-asian typefaces (a run naming only latin renders
    CJK in the theme's fallback face — the silent-substitution trap, OOXML
    edition)."""
    if not family:
        return
    from pptx.oxml.ns import qn
    run.font.name = family
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            from lxml import etree
            el = etree.SubElement(rPr, qn(tag))
        el.set('typeface', family)


def _add_shadow(shape, shadow: dict, theme: dict) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree
    spPr = shape._element.spPr
    effect = spPr.find(qn('a:effectLst'))
    if effect is None:
        effect = etree.SubElement(spPr, qn('a:effectLst'))
    shd = etree.SubElement(effect, qn('a:outerShdw'))
    blur = float(shadow.get('blur', 6))
    ox, oy = (shadow.get('offset') or [0, 0])[:2]
    import math
    dist = math.hypot(ox, oy)
    direction = int(math.degrees(math.atan2(oy, ox)) * 60000) % 21600000
    shd.set('blurRad', str(_emu(blur)))
    shd.set('dist', str(_emu(dist)))
    shd.set('dir', str(direction))
    color = (shadow.get('color') or '#00000040')
    from lib.slides.pptd import resolve_color
    color = resolve_color(color, theme, '#00000040')
    alpha = 0.4
    if re.match(r'^#[0-9a-fA-F]{8}$', color):
        alpha = int(color[7:9], 16) / 255.0
        color = color[:7]
    srgb = etree.SubElement(shd, qn('a:srgbClr'))
    srgb.set('val', color.lstrip('#'))
    a = etree.SubElement(srgb, qn('a:alpha'))
    a.set('val', str(int(alpha * 100000)))


# ── Text ──────────────────────────────────────────────────

_ALIGN_MAP = {'left': 1, 'center': 2, 'right': 3, 'justify': 4}


def _fill_text_frame(tf, content: dict, theme: dict, deck: Deck) -> None:
    from pptx.util import Pt
    st = text_style(content, theme)
    paragraphs = parse_rich_text(str(content.get('text') or ''), theme)
    tf.word_wrap = bool(content.get('wrap', True))
    align = content.get('align') or ['left', 'top']
    halign, valign = (list(align) + ['left', 'top'])[:2]
    anchor = {'top': 1, 'middle': 3, 'bottom': 4}.get(valign, 1)
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = {1: MSO_ANCHOR.TOP, 3: MSO_ANCHOR.MIDDLE,
                          4: MSO_ANCHOR.BOTTOM}[anchor]
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    if not paragraphs:
        paragraphs = []
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pa = para.align or halign
        if pa in _ALIGN_MAP:
            from pptx.enum.text import PP_ALIGN
            p.alignment = {1: PP_ALIGN.LEFT, 2: PP_ALIGN.CENTER,
                           3: PP_ALIGN.RIGHT, 4: PP_ALIGN.JUSTIFY}[
                               _ALIGN_MAP[pa]]
        lh_px = para.line_height_px or st.get('lineHeightPx')
        lh = para.line_height or st.get('lineHeight')
        if lh_px:
            p.line_spacing = Pt(float(lh_px))
        elif lh and float(lh) != 1.0:
            p.line_spacing = float(lh)
        mt = para.margin_top or st.get('marginTop')
        if mt:
            p.space_before = Pt(float(mt))
        if para.list_kind:
            prefix = f'{para.list_index}. ' if para.list_kind == 'ol' else '• '
            if para.runs:
                para.runs[0].text = prefix + para.runs[0].text
        for run_data in para.runs:
            for piece in run_data.text.split('\n'):
                run = p.add_run()
                run.text = piece
                if piece is not run_data.text.split('\n')[-1]:
                    run.text = piece + '\v'   # vertical tab = soft break
                f = run.font
                f.size = Pt(float(run_data.font_size
                                  or st.get('fontSize') or 18))
                f.bold = bool(run_data.bold or st.get('bold'))
                f.italic = bool(run_data.italic or st.get('italic'))
                f.underline = bool(run_data.underline)
                color = run_data.color or st.get('color') or '#000000'
                rgb, alpha = _rgb(color)
                f.color.rgb = rgb
                _set_run_font(run, run_data.font_family
                              or st.get('fontFamily') or 'MiSans')
                ls = st.get('letterSpacing')
                if ls:
                    run._r.get_or_add_rPr().set(
                        'spc', str(int(float(ls) * 100)))
                if run_data.sup or run_data.sub:
                    run._r.get_or_add_rPr().set(
                        'baseline', '30000' if run_data.sup else '-25000')
                if run_data.link:
                    run.hyperlink.address = run_data.link


def _add_text(slide, el: dict, deck: Deck) -> None:
    from pptx.util import Emu
    x, y, w, h = [float(v) for v in el['bounds']]
    tb = slide.shapes.add_textbox(Emu(_emu(x)), Emu(_emu(y)),
                                  Emu(_emu(w)), Emu(_emu(h)))
    tb.name = str(el.get('elementId') or 'text')
    _fill_text_frame(tb.text_frame, el.get('content') or {}, deck.theme, deck)
    if el.get('rotation'):
        tb.rotation = float(el['rotation'])


# ── Shapes ────────────────────────────────────────────────

_SHAPE_MAP = {
    'rect': 'RECTANGLE', 'roundRect': 'ROUNDED_RECTANGLE', 'ellipse': 'OVAL',
    'triangle': 'ISOCELES_TRIANGLE', 'diamond': 'DIAMOND',
    'homePlate': 'PENTAGON', 'chevron': 'CHEVRON', 'donut': 'DONUT',
    'star5': 'STAR_5_POINT', 'rightArrow': 'RIGHT_ARROW',
    'leftArrow': 'LEFT_ARROW', 'upArrow': 'UP_ARROW',
    'downArrow': 'DOWN_ARROW', 'leftRightArrow': 'LEFT_RIGHT_ARROW',
    'pentagon': 'REGULAR_PENTAGON', 'hexagon': 'HEXAGON',
    'parallelogram': 'PARALLELOGRAM', 'trapezoid': 'TRAPEZOID',
    'cross': 'CROSS', 'heart': 'HEART', 'lightningBolt': 'LIGHTNING_BOLT',
    'cloud': 'CLOUD', 'bracketPair': 'DOUBLE_BRACKET',
    'bracePair': 'DOUBLE_BRACE', 'round1Rect': 'ROUND_1_RECTANGLE',
    'round2SameRect': 'ROUND_2_SAME_RECTANGLE',
    'wedgeRectCallout': 'RECTANGULAR_CALLOUT',
}


def _apply_fill(shape_or_fill, fill: dict, theme: dict) -> None:
    """Solid/gradient/image fill onto a shape (image via pattern is not an
    OOXML shape fill — use a cropped picture for those; here: first stop)."""
    f = shape_or_fill
    if not isinstance(fill, dict):
        f.background()
        return
    ftype = fill.get('type')
    if ftype == 'solid':
        rgb, alpha = _rgb(_resolve(fill.get('color'), theme))
        f.solid()
        f.fore_color.rgb = rgb
        _set_alpha_on_fill(f, alpha)
    elif ftype == 'gradient':
        f.solid()
        rgb, _a = _rgb(_resolve((fill.get('stops') or [{}])[0].get('color'),
                                theme))
        f.fore_color.rgb = rgb
        _gradient_fill_xml(f._xPr, fill, theme)
    elif ftype == 'image':
        f.solid()
        f.fore_color.rgb = _rgb('#888888')[0]
    else:
        f.background()


def _resolve(value, theme):
    from lib.slides.pptd import resolve_color
    return resolve_color(value, theme, '#000000')


def _apply_border(line, border: dict, theme: dict) -> None:
    from pptx.util import Pt
    if not isinstance(border, dict):
        line.fill.background()
        return
    rgb, alpha = _rgb(_resolve(border.get('color'), theme))
    line.color.rgb = rgb
    line.width = Pt(float(border.get('width', 1)))
    style = border.get('style')
    if style in ('dash', 'dot'):
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = line._get_or_add_ln()
        dash = etree.SubElement(ln, qn('a:prstDash'))
        dash.set('val', 'dash' if style == 'dash' else 'sysDot')


def _add_shape(slide, el: dict, deck: Deck) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu
    x, y, w, h = [float(v) for v in el['bounds']]
    name = str(el.get('shapeName') or 'rect')
    if name == 'custom':
        _add_custom_shape(slide, el, deck, x, y, w, h)
        return
    mso_name = _SHAPE_MAP.get(name)
    if mso_name is None:
        logger.info('[Slides→PPTX] shape %s → rect fallback', name)
        mso_name = 'RECTANGLE'
    mso = getattr(MSO_SHAPE, mso_name, MSO_SHAPE.RECTANGLE)
    shp = slide.shapes.add_shape(mso, Emu(_emu(x)), Emu(_emu(y)),
                                 Emu(_emu(w)), Emu(_emu(h)))
    shp.name = str(el.get('elementId') or 'shape')
    if el.get('adjustments') and mso_name not in ('rect', 'ellipse',
                                                  'diamond'):
        try:
            for i, adj in enumerate(el['adjustments']):
                shp.adjustments[i] = float(adj) / 100000.0
        except (IndexError, ValueError) as e:
            logger.debug('[Slides→PPTX] adjustments for %s: %s', name, e)
    _apply_fill(shp.fill, el.get('fill'), deck.theme)
    _apply_border(shp.line, el.get('border'), deck.theme)
    if el.get('shadow'):
        _add_shadow(shp, el['shadow'], deck.theme)
    if el.get('rotation'):
        shp.rotation = float(el['rotation'])
    shp.shadow.inherit = False
    if el.get('opacity') is not None and float(el['opacity']) < 1:
        # Whole-shape opacity: alpha on the solid fill is the v1 answer.
        fill = el.get('fill') or {}
        if fill.get('type') == 'solid':
            rgb, alpha = _rgb(_resolve(fill.get('color'), deck.theme))
            shp.fill.solid()
            shp.fill.fore_color.rgb = rgb
            _set_alpha_on_fill(shp.fill, alpha * float(el['opacity']))


def _add_custom_shape(slide, el, deck, x, y, w, h):
    """shapeName=custom: SVG path → freeform polyline (curves sampled)."""
    from lib.slides.render_html import _shape_path  # noqa: F401  (unused: custom uses its own path)
    path = str(el.get('path') or '')
    vb = el.get('viewBox') or [w, h]
    pts = _sample_svg_path(path, float(vb[0]), float(vb[1]))
    if not pts:
        return
    sx, sy = w / float(vb[0]), h / float(vb[1])
    _add_freeform(slide, [(px * sx + x, py * sy + y) for px, py in pts],
                  el, deck, close=True)


def _sample_svg_path(d: str, vbw: float, vbh: float, *, per_curve: int = 12):
    """Sample an SVG path (M/L/H/V/C/Q/A→polyline) into viewBox points."""
    import math
    tokens = re.findall(r'[MLHVCQAZ]|-?\d*\.?\d+', d or '')
    pts: list = []
    cur = (0.0, 0.0)
    start = None
    i = 0
    cmd = ''

    def num():
        nonlocal i
        v = float(tokens[i])
        i += 1
        return v

    while i < len(tokens):
        if re.match(r'[MLHVCQAZ]', tokens[i]):
            cmd = tokens[i]
            i += 1
            continue
        try:
            if cmd == 'M':
                cur = (num(), num())
                pts.append(cur)
                start = cur
                cmd = 'L'
            elif cmd == 'L':
                cur = (num(), num())
                pts.append(cur)
            elif cmd == 'H':
                cur = (num(), cur[1])
                pts.append(cur)
            elif cmd == 'V':
                cur = (cur[0], num())
                pts.append(cur)
            elif cmd == 'C':
                c1 = (num(), num())
                c2 = (num(), num())
                end = (num(), num())
                p0 = cur
                for k in range(1, per_curve + 1):
                    t = k / per_curve
                    mt = 1 - t
                    bx = (mt**3 * p0[0] + 3 * mt**2 * t * c1[0]
                          + 3 * mt * t**2 * c2[0] + t**3 * end[0])
                    by = (mt**3 * p0[1] + 3 * mt**2 * t * c1[1]
                          + 3 * mt * t**2 * c2[1] + t**3 * end[1])
                    pts.append((bx, by))
                cur = end
            elif cmd == 'Q':
                c = (num(), num())
                end = (num(), num())
                p0 = cur
                for k in range(1, per_curve + 1):
                    t = k / per_curve
                    mt = 1 - t
                    pts.append((mt**2 * p0[0] + 2 * mt * t * c[0]
                                + t**2 * end[0],
                                mt**2 * p0[1] + 2 * mt * t * c[1]
                                + t**2 * end[1]))
                cur = end
            elif cmd == 'A':
                # Arc → sample as a circle/ellipse segment.
                rx, ry = num(), num()
                num()                     # rotation
                large, sweep = num(), num()
                end = (num(), num())
                p0 = cur
                # Endpoint → center parametrisation (SVG spec F.6.5).
                x1p = (p0[0] - end[0]) / 2
                y1p = (p0[1] - end[1]) / 2
                lam = (x1p**2 / rx**2 + y1p**2 / ry**2) if rx and ry else 99
                if lam > 1:
                    s = math.sqrt(lam)
                    rx, ry = rx * s, ry * s
                num_ = (rx**2 * ry**2 - rx**2 * y1p**2 - ry**2 * x1p**2)
                den = (rx**2 * y1p**2 + ry**2 * x1p**2) or 1
                coef = math.sqrt(max(0.0, num_ / den))
                if large == sweep:
                    coef = -coef
                cxp, cyp = coef * rx * y1p / (ry or 1), -coef * ry * x1p / (rx or 1)
                cx, cy = cxp + (p0[0] + end[0]) / 2, cyp + (p0[1] + end[1]) / 2
                a1 = math.atan2((p0[1] - cy) / (ry or 1),
                                (p0[0] - cx) / (rx or 1))
                a2 = math.atan2((end[1] - cy) / (ry or 1),
                                (end[0] - cx) / (rx or 1))
                if sweep and a2 < a1:
                    a2 += 2 * math.pi
                if not sweep and a2 > a1:
                    a2 -= 2 * math.pi
                for k in range(1, per_curve + 1):
                    a = a1 + (a2 - a1) * k / per_curve
                    pts.append((cx + rx * math.cos(a),
                                cy + ry * math.sin(a)))
                cur = end
            elif cmd == 'Z':
                if start:
                    pts.append(start)
                break
            else:
                i += 1
        except (ValueError, IndexError) as e:
            logger.debug('[Slides→PPTX] svg path sampling stopped at %d: %s',
                         i, e)
            break
    return pts


def _add_freeform(slide, points_pt, el, deck, *, close: bool = False):
    from pptx.util import Emu
    if len(points_pt) < 2:
        return
    builder = slide.shapes.build_freeform(Emu(_emu(points_pt[0][0])),
                                          Emu(_emu(points_pt[0][1])))
    builder.add_line_segments([(Emu(_emu(px)), Emu(_emu(py)))
                               for px, py in points_pt[1:]], close=close)
    shp = builder.convert_to_shape()
    shp.name = str(el.get('elementId') or 'freeform')
    if close:
        _apply_fill(shp.fill, el.get('fill'), deck.theme)
    else:
        shp.fill.background()
    _apply_border(shp.line, el.get('border'), deck.theme)
    shp.shadow.inherit = False
    arrow = el.get('arrow')
    if arrow and any(arrow):
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = shp.line._get_or_add_ln()
        if arrow[0]:
            he = etree.SubElement(ln, qn('a:headEnd'))
            he.set('type', 'triangle')
        if len(arrow) > 1 and arrow[1]:
            te = etree.SubElement(ln, qn('a:tailEnd'))
            te.set('type', 'triangle')


def _add_line(slide, el: dict, deck: Deck) -> None:
    from lib.slides.render_html import _line_path  # noqa: F401
    x, y, w, h = [float(v) for v in el['bounds']]
    vb = el.get('viewBox') or [w, h]
    pts = []
    for tok in str(el.get('points') or '').split():
        try:
            px, py = tok.split(',')
            pts.append((float(px), float(py)))
        except ValueError as e:
            logger.debug('[Slides→PPTX] skipping unparseable point %r: %s',
                         tok, e)
            continue
    if len(pts) < 2:
        return
    if str(el.get('curve') or 'round') == 'smooth' and len(pts) >= 3:
        sampled: list = [pts[0]]
        for i in range(len(pts) - 1):
            p0 = pts[i - 1] if i > 0 else pts[i]
            p1, p2 = pts[i], pts[i + 1]
            p3 = pts[i + 2] if i + 2 < len(pts) else p2
            for k in range(1, 9):
                t = k / 8
                c1 = (p1[0] + (p2[0] - p0[0]) / 6, p1[1] + (p2[1] - p0[1]) / 6)
                c2 = (p2[0] - (p3[0] - p1[0]) / 6, p2[1] - (p3[1] - p1[1]) / 6)
                mt = 1 - t
                sampled.append((mt**3 * p1[0] + 3 * mt**2 * t * c1[0]
                                + 3 * mt * t**2 * c2[0] + t**3 * p2[0],
                                mt**3 * p1[1] + 3 * mt**2 * t * c1[1]
                                + 3 * mt * t**2 * c2[1] + t**3 * p2[1]))
        pts = sampled
    sx, sy = w / float(vb[0] or w), h / float(vb[1] or h)
    _add_freeform(slide, [(px * sx + x, py * sy + y) for px, py in pts],
                  el, deck, close=False)


# ── Images ────────────────────────────────────────────────

def _image_size(path_or_url: str) -> tuple:
    try:
        from PIL import Image
        with Image.open(path_or_url) as im:
            return im.size
    except Exception as e:
        logger.debug('[Slides→PPTX] image size probe failed for %s: %s',
                     path_or_url, e)
        return (0, 0)


def _add_image(slide, el: dict, deck: Deck) -> None:
    from pptx.util import Emu
    x, y, w, h = [float(v) for v in el['bounds']]
    src = resolve_media(deck, str(el.get('src') or ''))
    if src.startswith(('http://', 'https://')):
        src = _download_media(deck, src)
        if not src:
            return
    fit = (el.get('fit') or {}).get('mode', 'cover')
    crop = el.get('crop') or {}
    cl = float(crop.get('left', 0) or 0)
    cr = float(crop.get('right', 0) or 0)
    ct = float(crop.get('top', 0) or 0)
    cb = float(crop.get('bottom', 0) or 0)

    if fit == 'contain':
        iw, ih = _image_size(src)
        rw = iw * (1 - cl - cr)
        rh = ih * (1 - ct - cb)
        if iw and ih and rw and rh:
            scale = min(w / rw, h / rh)
            pw, ph = rw * scale, rh * scale
            px, py = x + (w - pw) / 2, y + (h - ph) / 2
            pic = slide.shapes.add_picture(src, Emu(_emu(px)), Emu(_emu(py)),
                                           Emu(_emu(pw)), Emu(_emu(ph)))
            pic.crop_left, pic.crop_right = cl, cr
            pic.crop_top, pic.crop_bottom = ct, cb
            pic.name = str(el.get('elementId') or 'image')
            return
    pic = slide.shapes.add_picture(src, Emu(_emu(x)), Emu(_emu(y)),
                                   Emu(_emu(w)), Emu(_emu(h)))
    pic.name = str(el.get('elementId') or 'image')
    if fit == 'cover':
        iw, ih = _image_size(src)
        if iw and ih and (1 - cl - cr) > 0 and (1 - ct - cb) > 0:
            region_w = iw * (1 - cl - cr)
            region_h = ih * (1 - ct - cb)
            target = w / h
            region = region_w / region_h
            if region > target:     # too wide → crop more horizontally
                keep_w_frac = region_h * target / iw
                extra = (1 - cl - cr) - keep_w_frac
                cl += extra / 2
                cr += extra / 2
            else:                   # too tall → crop more vertically
                keep_h_frac = region_w / target / ih
                extra = (1 - ct - cb) - keep_h_frac
                ct += extra / 2
                cb += extra / 2
    if fit != 'fill':
        pic.crop_left = max(0.0, min(0.99, cl))
        pic.crop_right = max(0.0, min(0.99, cr))
        pic.crop_top = max(0.0, min(0.99, ct))
        pic.crop_bottom = max(0.0, min(0.99, cb))
    cs = el.get('cropShape') or {}
    cs_name = cs.get('shapeName')
    if cs_name and cs_name != 'rect':
        _apply_crop_shape(pic, cs_name, cs.get('adjustments'))
    if el.get('rotation'):
        pic.rotation = float(el['rotation'])


def _apply_crop_shape(pic, name: str, adjustments) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree
    mso = {'roundRect': 'roundRect', 'ellipse': 'ellipse',
           'triangle': 'triangle', 'diamond': 'diamond'}.get(name)
    if not mso:
        return
    spPr = pic._element.spPr
    geom = spPr.find(qn('a:prstGeom'))
    if geom is None:
        geom = etree.SubElement(spPr, qn('a:prstGeom'))
    geom.set('prst', mso)
    av = geom.find(qn('a:avLst'))
    if av is None:
        av = etree.SubElement(geom, qn('a:avLst'))
    if adjustments:
        for i, adj in enumerate(adjustments):
            gd = etree.SubElement(av, qn('a:gd'))
            gd.set('name', f'adj{i + 1}')
            gd.set('fmla', f'val {int(adj)}')


def _download_media(deck: Deck, url: str) -> str:
    """Remote media → cached local file (export must be offline-reproducible)."""
    try:
        from lib.design_sys._store import store_bytes
        from lib.http_client import http_get
        resp = http_get(url, timeout=60)
        data = getattr(resp, 'content', b'') or b''
        if getattr(resp, 'status_code', 0) != 200 or not data:
            logger.warning('[Slides→PPTX] media fetch HTTP %s: %s',
                           getattr(resp, 'status_code', '?'), url)
            return ''
        ext = os.path.splitext(url.split('?')[0])[1].lower()
        if ext not in ('.png', '.jpg', '.jpeg', '.gif', '.svg', '.webp'):
            ext = '.png'
        return store_bytes(data, name=f'remote{ext}', subdir='media')
    except Exception as e:
        logger.warning('[Slides→PPTX] media fetch failed %s: %s', url, e)
        return ''


# ── Icons (rasterised) ────────────────────────────────────

def _rasterize_icon(name: str, color: str, *, px: int = 256) -> str:
    """Built-in SVG glyph → transparent PNG path ('' on failure)."""
    from lib.slides.render_html import _ICON_PATHS
    path = _ICON_PATHS.get(name)
    if not path:
        return ''
    safe = re.sub(r'[^a-z0-9]+', '-', name.lower() + '-' + color.lstrip('#'))
    from lib.design_sys._store import store_root
    out_dir = os.path.join(store_root(), 'icons')
    os.makedirs(out_dir, exist_ok=True)
    out = os.path.join(out_dir, f'{safe}.png')
    if os.path.isfile(out):
        return out
    try:
        from playwright.sync_api import sync_playwright
        try:
            import chromium_env
            chromium_env.ensure_chromium_env(os.environ)
        except Exception as e:
            logger.debug('[Slides→PPTX] chromium shim: %s', e)
        html = (f'<!doctype html><body style="margin:0">'
                f'<svg width="{px}" height="{px}" viewBox="0 0 512 512">'
                f'<path d="{path}" fill="{color}"/></svg></body>')
        with sync_playwright() as p:
            browser = p.chromium.launch()
            try:
                page = browser.new_page(viewport={'width': px, 'height': px})
                page.set_content(html)
                page.wait_for_timeout(120)
                page.screenshot(path=out, omit_background=True)
            finally:
                browser.close()
        return out
    except Exception as e:
        logger.warning('[Slides→PPTX] icon rasterise failed for %s: %s',
                       name, e)
        return ''


def _add_icon(slide, el: dict, deck: Deck) -> None:
    from pptx.util import Emu
    x, y, w, h = [float(v) for v in el['bounds']]
    fill = el.get('fill') or {}
    color = _resolve(fill.get('color'), deck.theme) \
        if isinstance(fill, dict) else '#111111'
    png = _rasterize_icon(str(el.get('iconName') or ''), color)
    if not png:
        return
    pic = slide.shapes.add_picture(png, Emu(_emu(x)), Emu(_emu(y)),
                                   Emu(_emu(w)), Emu(_emu(h)))
    pic.name = str(el.get('elementId') or 'icon')


# ── Tables ────────────────────────────────────────────────

def _add_table(slide, el: dict, deck: Deck) -> None:
    from pptx.util import Emu
    x, y, w, h = [float(v) for v in el['bounds']]
    rows = el.get('rows') or []
    n_rows = len(rows)
    n_cols = max((len(r) for r in rows), default=0)
    if not n_rows or not n_cols:
        return
    # Expand the sparse rows (merged positions omitted per spec) into a grid.
    grid = [[None] * n_cols for _ in range(n_rows)]
    covered = set()
    for ri, row in enumerate(rows):
        ci = 0
        for cell in row:
            while (ri, ci) in covered:
                ci += 1
            if ci >= n_cols:
                break
            grid[ri][ci] = cell
            rs = int(cell.get('rowSpan') or 1)
            cs = int(cell.get('colSpan') or 1)
            for dr in range(rs):
                for dc in range(cs):
                    if dr or dc:
                        covered.add((ri + dr, ci + dc))
            ci += cs

    gfx = slide.shapes.add_table(n_rows, n_cols, Emu(_emu(x)), Emu(_emu(y)),
                                 Emu(_emu(w)), Emu(_emu(h)))
    tbl = gfx.table
    cw = el.get('columnWidths') or [1.0 / n_cols] * n_cols
    rh = el.get('rowHeights') or [1.0 / n_rows] * n_rows
    for ci in range(n_cols):
        tbl.columns[ci].width = Emu(_emu(w * float(cw[ci])))
    for ri in range(n_rows):
        tbl.rows[ri].height = Emu(_emu(h * float(rh[ri])))

    tstyle = table_style(el.get('style'), deck.theme)
    from lib.slides.render_html import _cell_style
    for ri in range(n_rows):
        for ci in range(n_cols):
            if (ri, ci) in covered and grid[ri][ci] is None:
                continue
            cell_data = grid[ri][ci] or {}
            cell = tbl.cell(ri, ci)
            rs = int(cell_data.get('rowSpan') or 1)
            cs = int(cell_data.get('colSpan') or 1)
            if rs > 1 or cs > 1:
                try:
                    cell.merge(tbl.cell(ri + rs - 1, ci + cs - 1))
                except Exception as e:
                    logger.debug('[Slides→PPTX] merge (%d,%d): %s', ri, ci, e)
            st = _cell_style(cell_data, ri, ci, n_rows, n_cols, tstyle,
                             deck.theme)
            from lib.slides.pptd import cell_content
            cc = cell_content(cell_data)
            tf = cell.text_frame
            tf.word_wrap = True
            content = {'text': cc.get('text') or '',
                       'align': cc.get('align') or st.get('align')
                                or ['center', 'middle']}
            merged_style = dict(st)
            _fill_text_frame(tf, {**content, **{k: v for k, v in
                             merged_style.items() if k in (
                                 'color', 'fontSize', 'fontFamily', 'bold',
                                 'italic', 'lineHeight', 'letterSpacing')}},
                             deck.theme, deck)
            fill = st.get('fill')
            if isinstance(fill, dict) and fill.get('type') == 'solid':
                rgb, alpha = _rgb(_resolve(fill.get('color'), deck.theme))
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb
                _set_alpha_on_fill(cell.fill, alpha)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb('#FFFFFF00'[:7])[0]
                from pptx.oxml.ns import qn as _qn
                # transparent: noFill on tcPr
                tcPr = cell._tc.get_or_add_tcPr()
                from lxml import etree
                for tag in ('a:solidFill', 'a:noFill'):
                    old = tcPr.find(_qn(tag))
                    if old is not None:
                        tcPr.remove(old)
                etree.SubElement(tcPr, _qn('a:noFill'))
            _apply_cell_borders(cell, st.get('border'), deck.theme)


def _apply_cell_borders(cell, border, theme) -> None:
    """Border (uniform dict form) → tcPr ln elements. None = theme hairline."""
    from pptx.oxml.ns import qn
    from lxml import etree
    if border is None:
        border = {'style': 'solid', 'width': 1, 'color': '$hairline'}
    if not isinstance(border, dict):
        return
    tcPr = cell._tc.get_or_add_tcPr()
    color = _resolve(border.get('color'), theme)
    rgb, alpha = _rgb(color)
    for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        old = tcPr.find(qn(tag))
        if old is not None:
            tcPr.remove(old)
    for tag in ('a:lnB', 'a:lnT', 'a:lnR', 'a:lnL'):  # schema order reversed
        ln = etree.SubElement(tcPr, qn(tag))
        ln.set('w', str(_emu(float(border.get('width', 1)))))
        ln.set('cap', 'flat')
        fill = etree.SubElement(ln, qn('a:solidFill'))
        srgb = etree.SubElement(fill, qn('a:srgbClr'))
        srgb.set('val', str(rgb))
        if alpha < 0.999:
            a = etree.SubElement(srgb, qn('a:alpha'))
            a.set('val', str(int(alpha * 100000)))
        if border.get('style') in ('dash', 'dot'):
            dash = etree.SubElement(ln, qn('a:prstDash'))
            dash.set('val', 'dash' if border['style'] == 'dash' else 'sysDot')
    # lnL..lnB must be the FIRST children of tcPr.
    for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        el = tcPr.find(qn(tag))
        tcPr.remove(el)
        tcPr.insert(0, el)
    order = ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']
    lns = [(order.index(t), tcPr.find(qn(t))) for t in order]
    for _, el in lns:
        tcPr.remove(el)
    for _, el in lns:
        if el is not None:
            tcPr.append(el)
    # move them to the front in order
    for t in reversed(order):
        el = tcPr.find(qn(t))
        if el is not None:
            tcPr.remove(el)
            tcPr.insert(0, el)


# ── Charts (native OOXML via python-pptx) ─────────────────

_CHART_XL = {
    'column': 'COLUMN_CLUSTERED', 'bar': 'BAR_CLUSTERED',
    'line': 'LINE', 'pie': 'PIE',
}


def _add_chart(slide, el: dict, deck: Deck) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Emu, Pt
    x, y, w, h = [float(v) for v in el['bounds']]
    ctype = str(el.get('chartType') or 'column')
    xl = getattr(XL_CHART_TYPE, _CHART_XL.get(ctype, 'COLUMN_CLUSTERED'))
    data = el.get('data') or {}
    cats = [str(c) for c in (data.get('categories') or [])]
    series = [s for s in (data.get('series') or []) if isinstance(s, dict)]
    if not cats or not series:
        return
    opts = el.get('options') or {}
    chart_data = CategoryChartData()
    chart_data.categories = cats
    for s in series:
        chart_data.add_series(str(s.get('name') or ''),
                              [float(v) for v in (s.get('values') or [])])
    gfx = slide.shapes.add_chart(xl, Emu(_emu(x)), Emu(_emu(y)),
                                 Emu(_emu(w)), Emu(_emu(h)), chart_data)
    chart = gfx.chart
    chart.has_title = False
    # Series colors: per-series override, else the theme cycle.
    from lib.slides.render_html import _chart_palette
    pal = _chart_palette(deck.theme)
    for i, ser in enumerate(chart.series):
        color = series[i].get('color') if i < len(series) else None
        hexv = _resolve(color, deck.theme) if color else pal[i % len(pal)]
        rgb, _a = _rgb(hexv)
        try:
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = rgb
            ser.format.line.color.rgb = rgb
        except Exception as e:
            logger.debug('[Slides→PPTX] series %d color: %s', i, e)
    show_legend = bool(opts.get('legend', len(series) > 1))
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(max(9.0, h * 0.036))
    # De-default discipline: no heavy gridlines; hairline category axis only.
    try:
        if ctype != 'pie':
            value_axis = chart.value_axis
            value_axis.has_major_gridlines = False
            value_axis.format.line.color.rgb = _rgb(
                _resolve('$hairline', deck.theme))[0]
            cat_axis = chart.category_axis
            cat_axis.format.line.color.rgb = _rgb(
                _resolve('$hairline', deck.theme))[0]
            for axis in (value_axis, cat_axis):
                axis.tick_labels.font.size = Pt(max(9.0, h * 0.04))
    except Exception as e:
        logger.debug('[Slides→PPTX] axis styling: %s', e)
    if opts.get('dataLabels'):
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            plot.data_labels.font.size = Pt(max(9.0, h * 0.036))
        except Exception as e:
            logger.debug('[Slides→PPTX] data labels: %s', e)


# ── Slide assembly ────────────────────────────────────────

def _set_background(slide, page: Page, deck: Deck) -> None:
    bg = page.background or {'type': 'solid', 'color': '#FFFFFF'}
    if bg.get('type') == 'solid':
        rgb, _a = _rgb(_resolve(bg.get('color'), deck.theme))
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb
        return
    # gradient / image: a locked full-bleed rectangle as the first shape.
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
                                 Emu(_emu(deck.width)), Emu(_emu(deck.height)))
    shp.name = '__background__'
    shp.line.fill.background()
    shp.shadow.inherit = False
    if bg.get('type') == 'gradient':
        _apply_fill(shp.fill, bg, deck.theme)
    elif bg.get('type') == 'image':
        src = resolve_media(deck, str(bg.get('src') or ''))
        if src.startswith(('http://', 'https://')):
            src = _download_media(deck, src)
        if src:
            # Replace the rectangle with the picture semantics: crop-cover.
            el = {'bounds': [0, 0, deck.width, deck.height], 'src':
                  bg.get('src'), 'fit': bg.get('fit') or {'mode': 'cover'},
                  'crop': bg.get('crop') or {}, 'elementId': '__background__'}
            sp = shp._element
            sp.getparent().remove(sp)
            _add_image(slide, el, deck)
            return
        _apply_fill(shp.fill, {'type': 'solid', 'color': '#111111'},
                    deck.theme)


def export_pptx(deck: Deck, out_path: str, *, transition: str = 'fade',
                embed_fonts: bool = True) -> dict:
    """Write the deck to a PPTX. Returns a summary dict.

    ``embed_fonts`` writes the deck's used typefaces as fntdata parts so the
    file renders with the intended faces on machines that lack them. Only
    TrueType-outline faces embed (woff2 is decompressed via fonttools when
    available); everything else is skipped with a log — the file always
    names the real family names regardless, so a machine with the faces
    installed renders them natively either way.
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(_emu(deck.width))
    prs.slide_height = Emu(_emu(deck.height))
    blank = prs.slide_layouts[6]

    for pi, page in enumerate(deck.pages):
        slide = prs.slides.add_slide(blank)
        _set_background(slide, page, deck)
        for el in page.elements:
            if not isinstance(el, dict):
                continue
            etype = el.get('elementType')
            try:
                if etype == 'text':
                    _add_text(slide, el, deck)
                elif etype == 'shape':
                    _add_shape(slide, el, deck)
                elif etype == 'line':
                    _add_line(slide, el, deck)
                elif etype == 'image':
                    _add_image(slide, el, deck)
                elif etype == 'icon':
                    _add_icon(slide, el, deck)
                elif etype == 'table':
                    _add_table(slide, el, deck)
                elif etype == 'chart':
                    _add_chart(slide, el, deck)
            except Exception as e:
                logger.error('[Slides→PPTX] page %d element %s failed: %s',
                             pi + 1, el.get('elementId'), e, exc_info=True)
        if page.notes:
            slide.notes_slide.notes_text_frame.text = page.notes

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)
    patched = 0
    if transition == 'fade':
        patched = _patch_fade_transitions(out_path)
    embedded = 0
    if embed_fonts:
        try:
            embedded = _embed_fonts(out_path, deck)
        except Exception as e:
            # Embedding is a polish layer — a failure here must never cost
            # the deck itself (the file names real family names, so machines
            # with the faces installed still render them).
            logger.warning('[Slides→PPTX] font embedding failed, shipping '
                           'unembedded: %s', e, exc_info=True)
    summary = _verify(out_path)
    summary.update({'output': out_path, 'slides': len(deck.pages),
                    'fadeTransitions': patched, 'embeddedFonts': embedded})
    logger.info('[Slides→PPTX] %s: %d slides, %d fade transitions, %d bytes',
                out_path, len(deck.pages), patched, summary['bytes'])
    return summary


# ── Font embedding (fntdata) ──────────────────────────────

#: fsType bits that FORBID embedding (0x0002 = restricted license).
_FSTYPE_RESTRICTED = 0x0002


def _glyphs_to_quadratic(glyphs, *, max_err: float = 1.0) -> dict:
    from fontTools.pens.cu2quPen import Cu2QuPen
    from fontTools.pens.ttGlyphPen import TTGlyphPen
    out = {}
    for name in glyphs.keys():
        glyph = glyphs[name]
        pen = TTGlyphPen(glyphs)
        cu2qu = Cu2QuPen(pen, max_err, reverse_direction=True)
        glyph.draw(cu2qu)
        out[name] = pen.glyph()
    return out


def _otf_to_ttf(font) -> None:
    """In-place CFF→TrueType conversion (fonttools otf2ttf recipe).

    cu2qu approximates the cubic CFF outlines with quadratic splines at
    1.0 max_err — visually lossless for screen type, and it turns the
    flagship CFF faces (MiSans et al.) into embeddable glyf fonts.
    """
    from fontTools.ttLib import newTable
    assert font.sfntVersion == 'OTTO' and 'CFF ' in font
    glyph_order = font.getGlyphOrder()
    font['loca'] = newTable('loca')
    font['glyf'] = glyf = newTable('glyf')
    glyf.glyphOrder = glyph_order
    glyf.glyphs = _glyphs_to_quadratic(font.getGlyphSet())
    del font['CFF ']
    if 'VORG' in font:
        del font['VORG']
    glyf.compile(font)
    hmtx = font['hmtx']
    for name, glyph in glyf.glyphs.items():
        if hasattr(glyph, 'xMin'):
            hmtx[name] = (hmtx[name][0], glyph.xMin)
    post = font['post']
    post.formatType = 2.0
    post.extraNames = []
    post.mapping = {}
    post.glyphOrder = glyph_order
    maxp = font['maxp']
    maxp.tableVersion = 0x00010000
    maxp.maxZones = 1
    maxp.maxTwilightPoints = 0
    maxp.maxStorage = 0
    maxp.maxFunctionDefs = 0
    maxp.maxInstructionDefs = 0
    maxp.maxStackElements = 0
    maxp.maxSizeOfInstructions = 0
    maxp.maxComponentElements = max(
        (len(g.components) if hasattr(g, 'components') else 0)
        for g in glyf.glyphs.values())
    maxp.compile(font)
    font['head'].glyphDataFormat = 0
    font.sfntVersion = '\x00\x01\x00\x00'


def _converted_glyf_ttf(font_id: str, weight: int, data: bytes) -> bytes:
    """CFF source bytes → glyf TTF bytes, cached in the design store.

    A 29k-glyph CJK face converts once (tens of seconds) and the result is
    content-addressed forever after.
    """
    from fontTools.ttLib import TTFont
    from lib.design_sys._store import store_root
    cache_dir = os.path.join(store_root(), 'fonts', font_id)
    cache = os.path.join(cache_dir, f'w{weight}-glyf.ttf')
    if os.path.isfile(cache):
        with open(cache, 'rb') as f:
            return f.read()
    import io as _io
    font = TTFont(_io.BytesIO(data))
    _otf_to_ttf(font)
    buf = _io.BytesIO()
    font.save(buf)
    out = buf.getvalue()
    os.makedirs(cache_dir, exist_ok=True)
    tmp = cache + '.tmp'
    with open(tmp, 'wb') as f:
        f.write(out)
    os.replace(tmp, cache)
    logger.info('[Slides→PPTX] converted %s w%d CFF→glyf TTF (%d bytes, '
                'cached)', font_id, weight, len(out))
    return out


def _font_file_for_embedding(font_id: str, weight: int) -> tuple:
    """(bytes, ok) for one registry face/weight — TrueType outlines only.

    PowerPoint's embedded-font reader wants sfnt bytes with a glyf table:
    TTF passes through, woff2 is decompressed via fonttools (an optional
    dependency, auto-installed once when missing), CFF-only OTF is skipped
    (PowerPoint rejects CFF fntdata on several versions — a skip is honest,
    a corrupted file is not).
    """
    from lib.design_sys import fonts as _fonts
    path = _fonts.ensure_font(font_id, weight)
    if not path:
        return b'', False
    with open(path, 'rb') as f:
        data = f.read()
    if path.lower().endswith('.ttf'):
        return data, True
    try:
        from fontTools.ttLib import TTFont
    except ImportError:
        try:
            import subprocess
            import sys as _sys
            subprocess.run([_sys.executable, '-m', 'pip', 'install', '--user',
                            'fonttools', 'brotli'],
                           capture_output=True, timeout=300, check=False)
            from fontTools.ttLib import TTFont
        except Exception as e:
            logger.warning('[Slides→PPTX] fonttools unavailable, embedding '
                           'skipped for %s: %s', font_id, e)
            return b'', False
    try:
        import io as _io
        font = TTFont(_io.BytesIO(data))
        if int(font['OS/2'].fsType) & _FSTYPE_RESTRICTED:
            logger.info('[Slides→PPTX] %s fsType is restricted — skipped',
                        font_id)
            return b'', False
        if 'glyf' not in font:
            # CFF-outline face (MiSans et al.): convert to TrueType outlines
            # via cu2qu (visually lossless at 1.0 max_err; cached after the
            # one-time conversion).
            return _converted_glyf_ttf(font_id, weight, data), True
        font.flavor = None
        buf = _io.BytesIO()
        font.save(buf)
        return buf.getvalue(), True
    except Exception as e:
        logger.warning('[Slides→PPTX] font conversion failed for %s: %s',
                       font_id, e)
        return b'', False


def _embed_fonts(pptx_path: str, deck: Deck) -> int:
    """Embed the deck's used typefaces as fntdata parts. Returns count.

    OOXML wiring: fntdata content-type default + presentation rels +
    ``embedTrueTypeFonts`` attr + ``<p:embeddedFontLst>`` inserted directly
    after ``<p:notesSz>`` (CT_Presentation child order — PowerPoint ignores
    a misplaced list the same way it ignores a misplaced transition).
    """
    from lib.design_sys import fonts as _fonts
    from lib.slides.render_html import collect_families
    wanted = collect_families(deck)
    by_family = {f.family: f for f in _fonts.FONT_REGISTRY}
    embeds: list = []                    # [(family, {weight: bytes})]
    for fam in sorted(wanted):
        face = by_family.get(fam)
        if face is None:
            continue
        weights: dict = {}
        for src in face.sources:
            data, ok = _font_file_for_embedding(face.id, src.weight)
            if ok:
                weights[src.weight] = data
        if weights:
            embeds.append((fam, weights))
    if not embeds:
        return 0

    with zipfile.ZipFile(pptx_path, 'r') as z:
        items = {i.filename: z.read(i.filename)
                 for i in z.infolist()}
        order = [i for i in z.infolist()]

    ct = items['[Content_Types].xml'].decode('utf-8')
    if 'fntdata' not in ct:
        ct = ct.replace('</Types>',
                        '<Default Extension="fntdata" ContentType="'
                        'application/x-fontdata"/></Types>')
        items['[Content_Types].xml'] = ct.encode('utf-8')

    rels_name = 'ppt/_rels/presentation.xml.rels'
    rels = items[rels_name].decode('utf-8')
    import re as _re
    existing_ids = [int(m) for m in _re.findall(r'Id="rId(\d+)"', rels)]
    next_id = max(existing_ids or [0]) + 1
    pres_name = 'ppt/presentation.xml'
    pres = items[pres_name].decode('utf-8')

    font_entries = []
    new_rels = []
    for fi, (fam, weights) in enumerate(embeds, 1):
        # PowerPoint's embed model has exactly four slots (regular / bold /
        # italic / boldItalic). Pick the weight nearest 400 as regular and
        # the nearest 700 as bold — at most one each.
        ws = sorted(weights)
        regular = min(ws, key=lambda w: abs(w - 400))
        bolds = [w for w in ws if w != regular]
        picks = [('regular', regular)]
        if bolds:
            picks.append(('bold', min(bolds, key=lambda w: abs(w - 700))))
        slots = []
        for tag, weight in picks:
            rid = f'rId{next_id}'
            next_id += 1
            part = f'ppt/fonts/font{next_id}.fntdata'
            items[part] = weights[weight]
            new_rels.append(
                f'<Relationship Id="{rid}" Type="http://schemas.'
                f'openxmlformats.org/officeDocument/2006/relationships/font"'
                f' Target="fonts/font{next_id}.fntdata"/>')
            slots.append((tag, rid))
        slot_xml = ''.join(
            f'<p:{tag} r:id="{rid}"/>'
            for tag, rid in slots)
        fam_esc = fam.replace('&', '&amp;').replace('<', '&lt;') \
                     .replace('>', '&gt;').replace('"', '&quot;')
        font_entries.append(
            f'<p:embeddedFont><p:font typeface="{fam_esc}"/>'
            f'{slot_xml}</p:embeddedFont>')

    rels = rels.replace('</Relationships>', ''.join(new_rels)
                        + '</Relationships>')
    items[rels_name] = rels.encode('utf-8')

    if 'embedTrueTypeFonts' not in pres:
        pres = pres.replace('<p:presentation ',
                            '<p:presentation embedTrueTypeFonts="1" ', 1)
    lst = ('<p:embeddedFontLst>' + ''.join(font_entries)
           + '</p:embeddedFontLst>')
    if '<p:notesSz' in pres:
        pres = _re.sub(r'(<p:notesSz[^>]*/>)', r'\1' + lst, pres, count=1)
    else:
        pres = pres.replace('</p:presentation>',
                            lst + '</p:presentation>')
    items[pres_name] = pres.encode('utf-8')

    tmp = pptx_path + '.tmp'
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as z:
        written = set()
        for info in order:
            z.writestr(info, items[info.filename])
            written.add(info.filename)
        for name, data in items.items():
            if name not in written:
                z.writestr(name, data)
    os.replace(tmp, pptx_path)
    logger.info('[Slides→PPTX] embedded %d typeface(s): %s', len(embeds),
                ', '.join(f for f, _ in embeds))
    return len(embeds)


# ── Transitions + verification ────────────────────────────

_FADE_XML = '<p:transition spd="fast" advClick="1"><p:fade/></p:transition>'
_SLIDE_RE = re.compile(r'ppt/slides/slide\d+\.xml$')
_TRANSITION_RE = re.compile(r'<p:transition\b[^>]*(?:/>|>.*?</p:transition>)',
                            re.DOTALL)


def _patch_fade_transitions(pptx_path: str) -> int:
    """Insert one root-level fade transition per slide (CT_Slide order:
    after cSld/clrMapOvr, before timing/extLst — anywhere else Office
    silently ignores it)."""
    count = 0
    tmp = pptx_path + '.tmp'
    with zipfile.ZipFile(pptx_path, 'r') as src, \
            zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as dst:
        for info in src.infolist():
            data = src.read(info.filename)
            if _SLIDE_RE.match(info.filename):
                text = data.decode('utf-8')
                text = _TRANSITION_RE.sub('', text)
                anchor = None
                for tag in ('clrMapOvr', 'cSld'):
                    m = re.search(rf'<p:{tag}\b[^>]*(?:/>|>.*?</p:{tag}>)',
                                  text, re.DOTALL)
                    if m:
                        anchor = m
                if anchor is None:
                    raise ExportError(
                        f'{info.filename}: no cSld anchor for transition')
                text = (text[:anchor.end()] + _FADE_XML
                        + text[anchor.end():])
                data = text.encode('utf-8')
                count += 1
            dst.writestr(info, data)
    os.replace(tmp, pptx_path)
    return count


def _verify(pptx_path: str) -> dict:
    import xml.etree.ElementTree as ET
    with zipfile.ZipFile(pptx_path) as z:
        broken = z.testzip()
        if broken:
            raise ExportError(f'PPTX CRC failed at {broken}')
        slides = [n for n in z.namelist() if _SLIDE_RE.match(n)]
        if not slides:
            raise ExportError('PPTX contains no slides')
        if 'ppt/presentation.xml' not in z.namelist():
            raise ExportError('PPTX missing presentation.xml')
        for name in slides:
            root = ET.fromstring(z.read(name))
            order = [c.tag.rsplit('}', 1)[-1] for c in root]
            if 'transition' in order:
                ti = order.index('transition')
                if 'cSld' in order and order.index('cSld') > ti:
                    raise ExportError(f'{name}: transition before cSld')
    return {'bytes': os.path.getsize(pptx_path)}
