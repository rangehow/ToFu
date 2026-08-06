"""lib/slides/import_pptx.py — PPTX → PPTD import (template replication).

(docs/SLIDES_CAPABILITY_DESIGN.md §4.6 / P4). The user-facing flow: 上传一个
现有 PPT/PPTX 模板 → 迁出主题(配色/字体/页面尺寸)+ 逐页元素 → PPTD 工程,
之后 produce_slides 用这套视觉语言写新内容。

Fidelity posture (deliberate, documented):

  * FULL fidelity: page size, theme palette (theme1.xml accent/dk/lt ramp),
    major/minor fonts (latin + ea), text boxes (position, runs, sizes,
    colors, bold/italic, alignment), pictures (bytes → media/), autoshapes
    (preset name, fill, border), native tables.
  * NOT imported (logged, never silent): charts (their XML stays a chart
    only in the source), groups (flattened one level), SmartArt, media,
    freeform geometry (approximated by its bounding box + fill), animations.
  * Everything unmapped lands in ``page.notes`` so no content is LOST even
    when it is not re-layoutable.
"""

from __future__ import annotations

import os
import re

from lib.log import get_logger

logger = get_logger(__name__)

__all__ = ['import_pptx', 'ImportReport']

_EMU_PER_PT = 12700


class ImportReport(dict):
    """{'deck_dir', 'manifest', 'pages', 'theme', 'skipped': [...]}."""


def _pt(v) -> float:
    return round(float(v) / _EMU_PER_PT, 2)


def _hex(color_format) -> str:
    try:
        rgb = color_format.rgb
        if rgb is not None:
            return f'#{rgb}'
    except Exception as e:
        logger.debug('[Slides←PPTX] color read failed: %s', e)
    try:
        tc = color_format.theme_color
        if tc is not None:
            return f'$theme_{tc}'
    except Exception as e:
        logger.debug('[Slides←PPTX] theme color read failed: %s', e)
    return ''


def _extract_theme(prs) -> dict:
    """Slide size + palette + fonts from the package's theme part."""
    out = {'size': [_pt(prs.slide_width), _pt(prs.slide_height)],
           'colors': {}, 'fonts': {}}
    # Parse the theme XML directly for the color scheme + font scheme.
    try:
        from lxml import etree
        pkg = prs.part.package
        for part in pkg.iter_parts():
            if not part.partname.endswith('theme1.xml'):
                continue
            root = etree.fromstring(part.blob)
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            scheme = root.find('.//a:clrScheme', ns)
            if scheme is not None:
                for child in scheme:
                    tag = etree.QName(child).localname
                    srgb = child.find('a:srgbClr', ns)
                    sysc = child.find('a:sysClr', ns)
                    if srgb is not None:
                        out['colors'][tag] = f'#{srgb.get("val")}'
                    elif sysc is not None:
                        out['colors'][tag] = f'#{sysc.get("lastClr")}'
            font_scheme = root.find('.//a:fontScheme', ns)
            if font_scheme is not None:
                for which in ('majorFont', 'minorFont'):
                    node = font_scheme.find(f'a:{which}', ns)
                    if node is None:
                        continue
                    latin = node.find('a:latin', ns)
                    ea = node.find('a:ea', ns)
                    out['fonts'][which] = {
                        'latin': latin.get('typeface') if latin is not None else '',
                        'ea': ea.get('typeface') if ea is not None else '',
                    }
            break
    except Exception as e:
        logger.warning('[Slides←PPTX] theme extraction failed: %s', e)
    return out


def _text_element(shape) -> dict | None:
    if not getattr(shape, 'has_text_frame', False):
        return None
    tf = shape.text_frame
    paras = []
    for p in tf.paragraphs:
        runs = []
        for r in p.runs:
            if not r.text:
                continue
            style_bits = []
            if r.font.size:
                style_bits.append(f'font-size:{_pt(r.font.size)}px')
            color = _hex(r.font.color)
            if color and not color.startswith('$theme_'):
                style_bits.append(f'color:{color}')
            text = r.text
            if style_bits:
                text = (f'<span style="{";".join(style_bits)}">'
                        f'{text}</span>')
            if r.font.bold:
                text = f'<strong>{text}</strong>'
            if r.font.italic:
                text = f'<em>{text}</em>'
            runs.append(text)
        if runs:
            align = ''
            try:
                if p.alignment is not None:
                    align = {1: 'left', 2: 'center', 3: 'right',
                             4: 'justify'}.get(int(p.alignment), '')
            except Exception as e:
                logger.debug('[Slides←PPTX] alignment read: %s', e)
            style = f' style="text-align:{align}"' if align else ''
            paras.append(f'<p{style}>{"".join(runs)}</p>')
    if not paras:
        return None
    fs = 0.0
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.size:
                fs = max(fs, _pt(r.font.size))
    content = {'text': '\n'.join(paras)}
    if fs:
        content['fontSize'] = fs
    fam = ''
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.name:
                fam = r.font.name
                break
        if fam:
            break
    if fam:
        content['fontFamily'] = fam
    return {'elementId': shape.name or 'text', 'elementType': 'text',
            'bounds': [_pt(shape.left), _pt(shape.top),
                       _pt(shape.width), _pt(shape.height)],
            'content': content}


def _shape_element(shape) -> dict:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    el = {'elementId': shape.name or 'shape', 'elementType': 'shape',
          'bounds': [_pt(shape.left), _pt(shape.top),
                     _pt(shape.width), _pt(shape.height)],
          'shapeName': 'rect'}
    try:
        auto = shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        if auto and shape.auto_shape_type is not None:
            name = str(shape.auto_shape_type)
            m = re.search(r'(\w+) \(', name)
            token = m.group(1) if m else name
            remap = {'RECTANGLE': 'rect', 'ROUNDED_RECTANGLE': 'roundRect',
                     'OVAL': 'ellipse', 'ISOCELES_TRIANGLE': 'triangle',
                     'DIAMOND': 'diamond', 'CHEVRON': 'chevron',
                     'PENTAGON': 'homePlate', 'DONUT': 'donut',
                     'STAR_5_POINT': 'star5', 'RIGHT_ARROW': 'rightArrow',
                     'LEFT_ARROW': 'leftArrow', 'UP_ARROW': 'upArrow',
                     'DOWN_ARROW': 'downArrow', 'HEART': 'heart',
                     'CLOUD': 'cloud', 'LIGHTNING_BOLT': 'lightningBolt'}
            el['shapeName'] = remap.get(token, 'rect')
    except Exception as e:
        logger.debug('[Slides←PPTX] shape type read: %s', e)
    try:
        fill = shape.fill
        if fill.type is not None and int(fill.type) == 1:   # solid
            color = _hex(fill.fore_color)
            if color and not color.startswith('$theme_'):
                el['fill'] = {'type': 'solid', 'color': color}
    except Exception as e:
        logger.debug('[Slides←PPTX] fill read: %s', e)
    try:
        line = shape.line
        if line.fill.type is not None and int(line.fill.type) == 1:
            color = _hex(line.color)
            if color and not color.startswith('$theme_'):
                el['border'] = {'style': 'solid',
                                'width': _pt(line.width or 12700),
                                'color': color}
    except Exception as e:
        logger.debug('[Slides←PPTX] line read: %s', e)
    return el


def import_pptx(pptx_path: str, out_dir: str, *, title: str = '') -> dict:
    """Convert a PPTX into a PPTD project directory. Returns ImportReport."""
    import yaml
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)
    os.makedirs(os.path.join(out_dir, 'pages'), exist_ok=True)
    os.makedirs(os.path.join(out_dir, 'media'), exist_ok=True)
    theme = _extract_theme(prs)
    skipped: list = []
    page_files: list = []

    for si, slide in enumerate(prs.slides, 1):
        elements: list = []
        dropped: list = []
        for shape in slide.shapes:
            try:
                st = shape.shape_type
                if st == MSO_SHAPE_TYPE.PICTURE:
                    ext = os.path.splitext(shape.image.filename or '')[1] \
                        or '.png'
                    media_name = f'slide{si:02d}_{len(elements):02d}{ext}'
                    with open(os.path.join(out_dir, 'media', media_name),
                              'wb') as f:
                        f.write(shape.image.blob)
                    elements.append({
                        'elementId': shape.name or 'image',
                        'elementType': 'image',
                        'bounds': [_pt(shape.left), _pt(shape.top),
                                   _pt(shape.width), _pt(shape.height)],
                        'src': f'media/{media_name}',
                        'fit': {'mode': 'cover'}})
                elif st == MSO_SHAPE_TYPE.TABLE or getattr(
                        shape, 'has_table', False):
                    tbl = shape.table
                    rows = []
                    for r in tbl.rows:
                        rows.append([{'text': c.text} for c in r.cells])
                    n_cols = len(tbl.columns)
                    elements.append({
                        'elementId': shape.name or 'table',
                        'elementType': 'table',
                        'bounds': [_pt(shape.left), _pt(shape.top),
                                   _pt(shape.width), _pt(shape.height)],
                        'columnWidths': [round(1.0 / n_cols, 4)] * n_cols,
                        'rowHeights': [round(1.0 / len(tbl.rows), 4)]
                                      * len(tbl.rows),
                        'rows': rows})
                elif st in (MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.MEDIA,
                            MSO_SHAPE_TYPE.GROUP):
                    dropped.append(f'{shape.name}({st})')
                elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    text_el = _text_element(shape)
                    if text_el and text_el['content']['text'].strip('<p>/'):
                        elements.append(text_el)
                    else:
                        elements.append(_shape_element(shape))
                elif getattr(shape, 'has_text_frame', False):
                    text_el = _text_element(shape)
                    if text_el:
                        elements.append(text_el)
            except Exception as e:
                logger.warning('[Slides←PPTX] slide %d shape %s failed: %s',
                               si, getattr(shape, 'name', '?'), e)
                dropped.append(getattr(shape, 'name', '?'))
        if dropped:
            skipped.append({'slide': si, 'shapes': dropped})
        page = {
            'pageType': 'content',
            'background': {'type': 'solid', 'color': '#FFFFFF'},
            'elements': elements,
        }
        if dropped:
            page['notes'] = '未导入元素: ' + ', '.join(dropped)
        name = f'pages/{si:02d}.page'
        with open(os.path.join(out_dir, name), 'w', encoding='utf-8') as f:
            yaml.safe_dump(page, f, allow_unicode=True, sort_keys=False)
        page_files.append(name)

    # Theme tokens: map the OOXML ramp onto our six-token system.
    colors = theme.get('colors') or {}
    dk1 = colors.get('dk1', '#1A1A1A')
    lt1 = colors.get('lt1', '#FFFFFF')
    accent1 = colors.get('accent1', '#16283C')
    accent2 = colors.get('accent2', '#C0652B')
    fonts = theme.get('fonts') or {}
    body_fam = (fonts.get('minorFont') or {}).get('ea') or \
        (fonts.get('minorFont') or {}).get('latin') or 'MiSans'
    display_fam = (fonts.get('majorFont') or {}).get('ea') or \
        (fonts.get('majorFont') or {}).get('latin') or body_fam
    manifest = {
        'version': 'v2',
        'title': title or os.path.splitext(os.path.basename(pptx_path))[0],
        'size': [int(theme['size'][0]), int(theme['size'][1])],
        'theme': {
            'colors': {'bg': lt1, 'ink': dk1, 'primary': accent1,
                       'accent': accent2, 'muted': '#8A93A3',
                       'hairline': '#D8D5CE'},
            'textStyles': {
                'title': {'fontSize': 36, 'color': '$primary', 'bold': True,
                          'fontFamily': display_fam},
                'body': {'fontSize': 18, 'color': '$ink',
                         'fontFamily': body_fam, 'lineHeight': 1.5},
                'caption': {'fontSize': 12, 'color': '$muted'},
                'bignum': {'fontSize': 80, 'color': '$accent', 'bold': True},
            },
        },
        'pages': page_files,
    }
    manifest_path = os.path.join(out_dir, 'deck.pptd')
    with open(manifest_path, 'w', encoding='utf-8') as f:
        yaml.safe_dump(manifest, f, allow_unicode=True, sort_keys=False)
    report = ImportReport(deck_dir=out_dir, manifest=manifest_path,
                          pages=len(page_files), theme=manifest['theme'],
                          skipped=skipped)
    logger.info('[Slides←PPTX] imported %s → %d pages, %d skipped group(s)',
                pptx_path, len(page_files), len(skipped))
    return report
