#!/usr/bin/env python3
"""Generate an editable PPTX version of the Tofu five-layer architecture diagram.

Every block is a native PowerPoint shape (rounded-rectangle) with editable text
and fill colour — no rasterised image, so authors can freely reposition or
relabel in PowerPoint / Keynote / WPS / Google Slides / LibreOffice Impress.

Usage:
    python scripts/build_architecture_pptx.py             # writes docs/architecture.pptx
    python scripts/build_architecture_pptx.py --out X.pptx

Requires: python-pptx (``pip install python-pptx`` or conda equivalent).
"""

from __future__ import annotations

import argparse
import os
import sys

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ImportError:  # pragma: no cover
    sys.stderr.write(
        'python-pptx is required. Install with:\n'
        '    pip install python-pptx\n'
        '  or, inside a conda env:\n'
        '    conda install -c conda-forge python-pptx\n'
    )
    raise


# ── Palette (matches docs/architecture_en.html) ──────────────
PALETTE = {
    'entry': {
        'lab_fill': 'FEF1E0', 'lab_border': 'E2A663', 'lab_text': '7A3E06',
        'tile_fill': 'FFFBF3', 'tile_border': 'E2A663', 'tile_text': '0F172A',
    },
    'orch': {
        'lab_fill': 'FBF3C8', 'lab_border': 'C9A42C', 'lab_text': '6B4800',
        'tile_fill': 'FFFDF0', 'tile_border': 'C9A42C', 'tile_text': '0F172A',
    },
    'ctx': {
        'lab_fill': 'E1ECFB', 'lab_border': '6A93CC', 'lab_text': '1E3A8A',
        'tile_fill': 'F5F9FF', 'tile_border': '6A93CC', 'tile_text': '0F172A',
    },
    'tools': {
        'lab_fill': 'E6EAFC', 'lab_border': '7A86CC', 'lab_text': '312E81',
        'tile_fill': 'F6F7FF', 'tile_border': '7A86CC', 'tile_text': '0F172A',
    },
    'llm': {
        'lab_fill': 'ECDEF9', 'lab_border': 'A08DC6', 'lab_text': '4C1D95',
        'tile_fill': 'FAF5FF', 'tile_border': 'A08DC6', 'tile_text': '0F172A',
    },
}

MUTED = '64748B'
LINE_STRONG = '94A3B8'
BUS_FILL = 'F8FAFC'


# ── Layout: 16:9 slide, logical coordinates in inches ────────
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN_X = 0.55
TITLE_Y = 0.30

LABEL_W = 1.45
ROW_X = MARGIN_X
CELLS_X = MARGIN_X + LABEL_W + 0.20
CELLS_W = SLIDE_W - CELLS_X - MARGIN_X

# 5 rows; row 4 (Tools) has an extra bus strip on top, so we give it more height
ROW_HEIGHTS = {
    'entry': 0.95,
    'orch':  0.95,
    'ctx':   0.95,
    'tools': 1.40,   # includes bus strip
    'llm':   0.95,
}
ROW_GAP = 0.18
ROWS_Y_START = 1.10


def _hex_to_rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


def _add_rounded(slide, x, y, w, h, *, fill, border, border_w=1.25,
                 corner=0.08):
    """Add a rounded rectangle shape and return it."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    # corner radius adjustment (0.0 to 0.5)
    shp.adjustments[0] = corner
    shp.fill.solid()
    shp.fill.fore_color.rgb = _hex_to_rgb(fill)
    shp.line.color.rgb = _hex_to_rgb(border)
    shp.line.width = Pt(border_w)
    # kill default shadow
    shp.shadow.inherit = False
    return shp


def _set_text(shape, lines, *, font_size_primary=14, font_size_secondary=10,
              primary_bold=True, secondary_italic=True, text_color='0F172A',
              secondary_color=MUTED, anchor_middle=True, align_center=True):
    """Populate a shape's text frame with up to two tiers of text.

    lines: list of (text, tier) where tier ∈ {'p', 's', 'tag'}.
      'p'    → primary / bold name
      's'    → secondary / italic tag
      'tag'  → small uppercase badge (layer label)
    """
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(60000)
    tf.margin_top = tf.margin_bottom = Emu(40000)
    if anchor_middle:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Reset: one paragraph already exists
    first = True
    for (text, tier) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.CENTER if align_center else PP_ALIGN.LEFT
        # Remove any prior runs
        for r in list(p.runs):
            r.text = ''
        run = p.add_run()
        run.text = text
        f = run.font
        if tier == 'p':
            f.bold = primary_bold
            f.size = Pt(font_size_primary)
            f.color.rgb = _hex_to_rgb(text_color)
            f.name = 'Calibri'
        elif tier == 's':
            f.italic = secondary_italic
            f.size = Pt(font_size_secondary)
            f.color.rgb = _hex_to_rgb(secondary_color)
            f.name = 'Cambria'
        elif tier == 'tag':
            f.bold = False
            f.size = Pt(8)
            f.color.rgb = _hex_to_rgb(secondary_color)
            f.name = 'Calibri'
            # letter-spacing faked with uppercase
            run.text = text.upper()


def _add_label(slide, row_key, y, h, label_lines):
    p = PALETTE[row_key]
    shp = _add_rounded(slide, ROW_X, y, LABEL_W, h,
                       fill=p['lab_fill'], border=p['lab_border'],
                       border_w=1.4, corner=0.18)
    _set_text(shp, label_lines,
              font_size_primary=14, font_size_secondary=8,
              text_color=p['lab_text'])
    return shp


def _add_tile(slide, row_key, x, y, w, h, name, tag):
    p = PALETTE[row_key]
    shp = _add_rounded(slide, x, y, w, h,
                       fill=p['tile_fill'], border=p['tile_border'],
                       border_w=1.2, corner=0.12)
    _set_text(shp,
              [(name, 'p'), (tag, 's')],
              font_size_primary=13, font_size_secondary=10,
              text_color=p['tile_text'])
    return shp


def _arrange_cells(row_x, row_w, weights):
    """Given total width and a list of column weights summing to 12,
    return (x, width) for each cell with a small gap between them."""
    total_units = sum(weights)
    gap = 0.12
    total_gap = gap * (len(weights) - 1)
    unit_w = (row_w - total_gap) / total_units
    positions = []
    cursor = row_x
    for i, w in enumerate(weights):
        wi = unit_w * w
        positions.append((cursor, wi))
        cursor += wi + gap
    return positions


def build(out_path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── Title ──
    tb = slide.shapes.add_textbox(
        Inches(MARGIN_X), Inches(TITLE_Y),
        Inches(SLIDE_W - 2 * MARGIN_X), Inches(0.55),
    )
    _set_text(tb,
              [('Tofu: A Five-Layer Agent Harness', 'p'),
               ('Figure 1. System overview.', 's')],
              font_size_primary=22, font_size_secondary=12,
              text_color='0F172A', secondary_color=MUTED)

    # ── Rows ──
    y = ROWS_Y_START

    # Layer 1 — Entry
    h = ROW_HEIGHTS['entry']
    _add_label(slide, 'entry', y, h,
               [('Entry', 'p'), ('LAYER 1', 'tag')])
    cells = _arrange_cells(CELLS_X, CELLS_W, [4, 4, 4])
    entry_items = [
        ('Web UI',        'browser front-end'),
        ('Group Bot',     'chat-app collaboration'),
        ('Headless API',  'backend for third-party CLIs'),
    ]
    for (cx, cw), (nm, tag) in zip(cells, entry_items):
        _add_tile(slide, 'entry', cx, y, cw, h, nm, tag)
    y += h + ROW_GAP

    # Layer 2 — Orchestration
    h = ROW_HEIGHTS['orch']
    _add_label(slide, 'orch', y, h,
               [('Orchestration', 'p'), ('LAYER 2', 'tag')])
    cells = _arrange_cells(CELLS_X, CELLS_W, [4, 4, 4])
    orch_items = [
        ('ReAct Loop',             'single agent'),
        ('Planner–Worker–Critic',  'three-role review'),
        ('Swarm',                  'parallel specialist DAG'),
    ]
    for (cx, cw), (nm, tag) in zip(cells, orch_items):
        _add_tile(slide, 'orch', cx, y, cw, h, nm, tag)
    y += h + ROW_GAP

    # Layer 3 — Context Engineering
    h = ROW_HEIGHTS['ctx']
    _add_label(slide, 'ctx', y, h,
               [('Context Engineering', 'p'), ('LAYER 3', 'tag')])
    cells = _arrange_cells(CELLS_X, CELLS_W, [3, 3, 3, 3])
    ctx_items = [
        ('Prompt Assembly', 'system message'),
        ('Compaction',      'three-layer compression'),
        ('Memory',          'cross-session knowledge'),
        ('Message Builder', 'DB to API messages'),
    ]
    for (cx, cw), (nm, tag) in zip(cells, ctx_items):
        _add_tile(slide, 'ctx', cx, y, cw, h, nm, tag)
    y += h + ROW_GAP

    # Layer 4 — Tools & Extensions (with bus strip on top)
    h_full = ROW_HEIGHTS['tools']
    bus_h = 0.40
    tiles_h = h_full - bus_h - 0.10
    _add_label(slide, 'tools', y, h_full,
               [('Tools & Extensions', 'p'), ('LAYER 4', 'tag')])
    # bus strip
    bus = _add_rounded(slide, CELLS_X, y, CELLS_W, bus_h,
                       fill=BUS_FILL, border=LINE_STRONG,
                       border_w=1.1, corner=0.18)
    bus.line.dash_style = 7  # dash (may render solid in some viewers; acceptable)
    _set_text(bus,
              [('Tool-execution bus   dispatch  →  executor  →  handlers/', 'p')],
              font_size_primary=12, primary_bold=True,
              text_color='0F172A')
    # four tool tiles underneath
    tiles_y = y + bus_h + 0.10
    cells = _arrange_cells(CELLS_X, CELLS_W, [3, 3, 3, 3])
    tool_items = [
        ('Project Co-pilot', 'code and shell'),
        ('Search & Fetch',   'web information'),
        ('Multi-modal I/O',  'image, PDF, browser'),
        ('Human & MCP',      'open extensions'),
    ]
    for (cx, cw), (nm, tag) in zip(cells, tool_items):
        _add_tile(slide, 'tools', cx, tiles_y, cw, tiles_h, nm, tag)
    y += h_full + ROW_GAP

    # Layer 5 — LLM Dispatch
    h = ROW_HEIGHTS['llm']
    _add_label(slide, 'llm', y, h,
               [('LLM Dispatch', 'p'), ('LAYER 5', 'tag')])
    cells = _arrange_cells(CELLS_X, CELLS_W, [3, 3, 3, 3])
    llm_items = [
        ('Slot Pool',        'key × model'),
        ('Unified API',      'fallback and retry'),
        ('Discovery',        'capability registry'),
        ('Streaming Client', 'vendor-agnostic SSE'),
    ]
    for (cx, cw), (nm, tag) in zip(cells, llm_items):
        _add_tile(slide, 'llm', cx, y, cw, h, nm, tag)
    y += h + ROW_GAP

    # Save
    prs.save(out_path)


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument('--out', default='docs/architecture.pptx',
                    help='Output .pptx path (default: docs/architecture.pptx)')
    args = ap.parse_args()

    out_path = os.path.abspath(args.out)
    os.makedirs(os.path.dirname(out_path) or '.', exist_ok=True)
    build(out_path)
    print(f'Wrote {out_path}')


if __name__ == '__main__':
    main()
