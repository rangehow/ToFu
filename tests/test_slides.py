"""tests/test_slides.py — slides capability contracts (P3).

Pins: PPTD parse/validate (the zero-LLM gate the author loop repairs
against), theme resolution, the rich-text parser, the HTML renderer's
determinism, the PPTX exporter's structural guarantees (real shapes, fade
order, CRC), the recipe's degrade discipline (a bad page never kills the
deck), and the runtime/starter kind contract that keeps /api/v1/tasks/*
able to poll what /start launches.
"""

from __future__ import annotations

import os
import sys

import pytest

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from lib.slides import pptd  # noqa: E402
from lib.slides.pptd import Deck, Page, parse_deck, validate_deck  # noqa: E402

pytestmark = pytest.mark.unit


# ── Fixtures ──────────────────────────────────────────────

def _write_deck(tmp_path, pages_yaml: list, *, theme=None, size=(1280, 720)):
    """Write a minimal deck dir; return manifest path."""
    import yaml
    deck_dir = tmp_path / 'deck'
    (deck_dir / 'pages').mkdir(parents=True)
    names = []
    for i, text in enumerate(pages_yaml, 1):
        name = f'pages/{i:02d}.page'
        (deck_dir / name).write_text(text, encoding='utf-8')
        names.append(name)
    manifest = {'version': 'v2', 'title': '测试 deck', 'size': list(size),
                'theme': theme or {
                    'colors': {'bg': '#F7F7F5', 'ink': '#1B2430',
                               'primary': '#16283C', 'accent': '#C0652B',
                               'muted': '#6B7280', 'hairline': '#D8D5CE'},
                    'textStyles': {'title': {'fontSize': 40,
                                             'color': '$primary'},
                                   'body': {'fontSize': 18, 'color': '$ink'}},
                },
                'pages': names}
    (deck_dir / 'deck.pptd').write_text(
        yaml.safe_dump(manifest, allow_unicode=True), encoding='utf-8')
    return str(deck_dir / 'deck.pptd')


_COVER = '''pageType: cover
background: {type: solid, color: "$bg"}
elements:
  - elementId: title
    elementType: text
    bounds: [72, 200, 1136, 120]
    content:
      style: "$title"
      align: [left, middle]
      text: |
        <p><strong>一寸万象</strong></p>
  - elementId: rule
    elementType: shape
    bounds: [72, 340, 64, 6]
    shapeName: rect
    fill: {type: solid, color: "$accent"}
'''

_TABLE_PAGE = '''pageType: content
background: {type: solid, color: "$bg"}
elements:
  - elementId: t
    elementType: table
    bounds: [80, 120, 1120, 300]
    columnWidths: [0.5, 0.5]
    rowHeights: [0.34, 0.33, 0.33]
    rows:
      - - text: "指标"
        - text: "2025"
      - - text: "营收"
        - text: "96.3"
      - - text: "利润"
        - text: "15.8"
'''


# ── parse / validate ──────────────────────────────────────

class TestParseValidate:
    def test_round_trip_clean(self, tmp_path):
        deck = parse_deck(_write_deck(tmp_path, [_COVER, _TABLE_PAGE]))
        assert deck.title == '测试 deck'
        assert deck.size == (1280, 720)
        assert len(deck.pages) == 2
        assert validate_deck(deck) == []

    def test_manifest_dir_resolution(self, tmp_path):
        path = _write_deck(tmp_path, [_COVER])
        deck = parse_deck(os.path.dirname(path))
        assert len(deck.pages) == 1

    def test_bad_version_rejected(self, tmp_path):
        path = _write_deck(tmp_path, [_COVER])
        text = open(path).read().replace('v2', 'v9')
        open(path, 'w').write(text)
        with pytest.raises(pptd.PPTDError):
            parse_deck(path)

    def test_path_escape_rejected(self, tmp_path):
        deck_dir = tmp_path / 'deck'
        (deck_dir / 'pages').mkdir(parents=True)
        (deck_dir / 'deck.pptd').write_text(
            'version: v2\nsize: [1280, 720]\npages: ["../escape.page"]\n',
            encoding='utf-8')
        with pytest.raises(pptd.PPTDError):
            parse_deck(str(deck_dir / 'deck.pptd'))

    def test_validator_catches_real_defects(self, tmp_path):
        bad = '''pageType: content
background: {type: solid, color: "$bg"}
elements:
  - elementId: a
    elementType: text
    bounds: [10, 10, 100, 50]
    content: {style: "$nosuch", text: "x"}
  - elementId: a
    elementType: text
    bounds: [10, 10, 100, 50]
    content: {text: "dup id"}
  - elementId: b
    elementType: shape
    bounds: [10, 10, 100, 50]
    shapeName: nonExistentShape
  - elementId: c
    elementType: table
    bounds: [10, 10, 100, 50]
    columnWidths: [0.5, 0.6]
    rowHeights: [1.0]
    rows: [[{text: "x"}, {text: "y"}]]
'''
        deck = parse_deck(_write_deck(tmp_path, [bad]))
        findings = validate_deck(deck)
        blob = '\n'.join(findings)
        assert 'unknown textStyle token' in blob
        assert 'duplicate elementId' in blob
        assert 'unsupported shape' in blob
        assert 'columnWidths must sum to 1' in blob

    def test_theme_resolution(self, tmp_path):
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        assert pptd.resolve_color('$primary', deck.theme) == '#16283C'
        assert pptd.resolve_color('#AABBCCDD', deck.theme) == '#AABBCCDD'
        assert pptd.resolve_color('$nope', deck.theme, 'x') == 'x'
        st = pptd.text_style({'style': '$title'}, deck.theme)
        assert st['fontSize'] == 40 and st['color'] == '#16283C'
        st2 = pptd.text_style({'style': '$title', 'fontSize': 22},
                              deck.theme)
        assert st2['fontSize'] == 22      # inline wins over the theme style


# ── rich text ─────────────────────────────────────────────

class TestRichText:
    def test_plain_shorthand(self):
        from lib.slides.richtext import parse_rich_text
        paras = parse_rich_text('第一行\n第二行', {})
        assert len(paras) == 2
        assert paras[0].runs[0].text == '第一行'

    def test_marks_and_inline_styles(self):
        from lib.slides.richtext import parse_rich_text
        paras = parse_rich_text(
            '<p style="text-align:center"><strong>重</strong>'
            '<span style="color:$primary;font-size:24px">点</span></p>',
            {'colors': {'primary': '#123456'}})
        runs = paras[0].runs
        assert paras[0].align == 'center'
        assert runs[0].bold and runs[0].text == '重'
        assert runs[1].color == '#123456' and runs[1].font_size == 24

    def test_lists(self):
        from lib.slides.richtext import parse_rich_text
        paras = parse_rich_text('<ul><li>甲</li><li>乙</li></ul>', {})
        assert len(paras) == 2 and paras[0].list_kind == 'ul'

    def test_malformed_degrades_to_text(self):
        from lib.slides.richtext import parse_rich_text
        paras = parse_rich_text('<p><strong>未闭合', {})
        assert paras and '未闭合' in paras[0].runs[0].text


# ── HTML renderer ─────────────────────────────────────────

class TestRenderHtml:
    def test_page_html_deterministic_and_themed(self, tmp_path):
        from lib.slides.render_html import render_page_html
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        a = render_page_html(deck, deck.pages[0])
        b = render_page_html(deck, deck.pages[0])
        assert a == b                       # deterministic
        assert '#F7F7F5' in a               # $bg resolved
        assert '#16283C' in a               # $primary resolved
        assert '1280px' in a and '720px' in a
        assert '一寸万象' in a

    def test_gradient_angle_mapping(self):
        from lib.slides.render_html import _gradient_css
        css = _gradient_css({'type': 'gradient', 'gradientType': 'linear',
                             'angle': 90,
                             'stops': [{'position': 0, 'color': '#000000'},
                                       {'position': 1, 'color': '#FFFFFF'}]},
                            {})
        assert '180.0deg' in css            # PPTD 90 (top→bottom) = CSS 180

    def test_table_and_richtext_render(self, tmp_path):
        from lib.slides.render_html import render_page_html
        deck = parse_deck(_write_deck(tmp_path, [_TABLE_PAGE]))
        html = render_page_html(deck, deck.pages[0])
        assert '<table' in html and '指标' in html and 'colspan' not in html


# ── PPTX exporter ─────────────────────────────────────────

class TestExportPptx:
    def test_export_structure(self, tmp_path):
        pytest.importorskip('pptx')
        from lib.slides.export_pptx import export_pptx
        deck = parse_deck(_write_deck(tmp_path, [_COVER, _TABLE_PAGE]))
        out = str(tmp_path / 'out.pptx')
        summary = export_pptx(deck, out)
        assert summary['slides'] == 2
        assert summary['fadeTransitions'] == 2
        assert summary['bytes'] > 4096

        from pptx import Presentation
        prs = Presentation(out)
        assert len(prs.slides) == 2
        assert prs.slide_width == 1280 * 12700
        s1, s2 = prs.slides
        texts = []
        for sh in s1.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
        assert any('一寸万象' in t for t in texts)
        assert any(sh.has_table for sh in s2.shapes)

    def test_fade_transition_order(self, tmp_path):
        pytest.importorskip('pptx')
        import re
        import zipfile
        from lib.slides.export_pptx import export_pptx
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        out = str(tmp_path / 'out.pptx')
        export_pptx(deck, out)
        with zipfile.ZipFile(out) as z:
            xml = z.read('ppt/slides/slide1.xml').decode()
        i_csld = xml.index('</p:cSld>')
        i_trans = xml.index('<p:transition')
        assert i_trans > i_csld            # CT_Slide order
        assert xml.count('<p:fade/>') == 1

    def test_notes_and_font_ea(self, tmp_path):
        pytest.importorskip('pptx')
        import zipfile
        from lib.slides.export_pptx import export_pptx
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        out = str(tmp_path / 'out.pptx')
        export_pptx(deck, out)
        with zipfile.ZipFile(out) as z:
            xml = z.read('ppt/slides/slide1.xml').decode()
        assert '<a:ea typeface=' in xml     # CJK font on the ea slot


# ── Recipe (mocked LLM) ───────────────────────────────────

class TestRecipe:
    def _fake_llm(self, content):
        def _call(messages, **kw):
            return content, {'total_tokens': 10}
        return _call

    def test_full_graph_and_resume(self, tmp_path, monkeypatch):
        from lib.slides import recipe
        outline_json = ('{"title": "T", "scenario": "tech-engineering",'
                        ' "pages": ['
                        ' {"pageType": "cover", "purpose": "p", '
                        '  "key_message": "标题判断", "layout_hint": "", '
                        '  "content_notes": "n"},'
                        ' {"pageType": "content", "purpose": "p2", '
                        '  "key_message": "第二页判断", "layout_hint": "", '
                        '  "content_notes": "n2"},'
                        ' {"pageType": "final", "purpose": "p3", '
                        '  "key_message": "结论", "layout_hint": "", '
                        '  "content_notes": "n3"}]}')
        calls = {'n': 0}

        def _llm(messages, **kw):
            calls['n'] += 1
            text = messages[0]['content']
            if '大纲' in text or 'outline' in text:
                return outline_json, {}
            # page author path: return a minimal valid page
            return ('pageType: content\n'
                    'background: {type: solid, color: "$bg"}\n'
                    'elements:\n'
                    '  - elementId: t\n'
                    '    elementType: text\n'
                    '    bounds: [72, 72, 1136, 120]\n'
                    '    content:\n'
                    '      style: "$title"\n'
                    '      text: |\n'
                    '        <p>标题</p>\n'), {}
        monkeypatch.setattr(recipe, '_llm_chat', _llm)
        import lib.slides.author as author
        monkeypatch.setattr(author, '_llm',
                            lambda messages, **kw: _llm(messages, **kw))
        # render + qa mocked: no browser, no VLM in the unit lane
        import lib.slides.recipe as r
        monkeypatch.setattr(r, '_run_render',
                            lambda ctx: {'previews': [], 'failed': []})
        import lib.design_sys.visual_qa as vqa
        monkeypatch.setattr(vqa, 'visual_qa_available',
                            lambda: (False, 'test'))

        out = recipe.build_deck_from_topic('测试主题', str(tmp_path / 'job'))
        assert os.path.isfile(out['pptx_path'])
        assert out['pages'] == 3
        assert out['authored_pages'] == 3
        assert out['theme_id'] == 'paper-engineer'

        # Resume: a second run must skip every finished stage (checkpoint).
        n_before = calls['n']
        out2 = recipe.build_deck_from_topic('测试主题', str(tmp_path / 'job'))
        assert calls['n'] == n_before     # zero LLM calls on full resume
        assert out2['pptx_path'] == out['pptx_path']

    def test_bad_page_degrades_not_fails(self, tmp_path, monkeypatch):
        from lib.slides import recipe
        import lib.slides.author as author
        outline_json = ('{"title": "T", "scenario": "business-plan",'
                        ' "pages": ['
                        ' {"pageType": "cover", "key_message": "甲"},'
                        ' {"pageType": "content", "key_message": "乙"},'
                        ' {"pageType": "final", "key_message": "丙"}]}')

        def _llm(messages, **kw):
            text = messages[0]['content']
            if '大纲' in text:
                return outline_json, {}
            return '这不是 YAML: [{{{', {}
        monkeypatch.setattr(recipe, '_llm_chat', _llm)
        monkeypatch.setattr(author, '_llm',
                            lambda messages, **kw: _llm(messages, **kw))
        monkeypatch.setattr(recipe, '_run_render',
                            lambda ctx: {'previews': [], 'failed': []})
        import lib.design_sys.visual_qa as vqa
        monkeypatch.setattr(vqa, 'visual_qa_available',
                            lambda: (False, 'test'))
        out = recipe.build_deck_from_topic('x', str(tmp_path / 'job'))
        assert os.path.isfile(out['pptx_path'])
        assert out['authored_pages'] == 0   # all fell back — deck survived
        assert out['pages'] == 3


# ── Runtime / starter contract ────────────────────────────

class TestRuntimeContract:
    def test_kind_matches_starter(self):
        """The kind /start dispatches on must be the kind the runtime
        registers — a mismatch means a job you can start but never poll."""
        from lib.slides.runtime import _slides_runtime
        assert _slides_runtime.kind == 'slides-deck'

    def test_fallback_page_validates(self, tmp_path):
        from lib.slides.author import fallback_page
        from lib.design_sys.themes import get_theme
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        theme = get_theme('paper-engineer')
        yaml_text = fallback_page(deck, {'key_message': '完整判断句'},
                                  theme=theme)
        page_path = os.path.join(deck.root, 'pages', '99_fb.page')
        open(page_path, 'w').write(yaml_text)
        deck2 = parse_deck(deck.manifest_path)
        deck2.pages.append(pptd.Page(path='pages/99_fb.page',
                                     elements=__import__('yaml').safe_load(
                                         yaml_text)['elements'],
                                     background={'type': 'solid',
                                                 'color': '$bg'}))
        assert validate_deck(deck2) == []

    def test_tool_schema_and_family(self):
        from lib.tools.produce import (PRODUCE_SLIDES_TOOL,
                                       PRODUCE_SLIDES_TOOL_NAME,
                                       PRODUCE_TOOL_NAMES)
        assert PRODUCE_SLIDES_TOOL_NAME in PRODUCE_TOOL_NAMES
        assert PRODUCE_SLIDES_TOOL['function']['name'] == 'produce_slides'
        assert 'topic' in PRODUCE_SLIDES_TOOL['function']['parameters'][
            'required']


# ── P4: native chart / font embedding / import round-trip ──

_CHART_PAGE = '''pageType: content
background: {type: solid, color: "$bg"}
elements:
  - elementId: c
    elementType: chart
    bounds: [80, 120, 600, 360]
    chartType: bar
    data:
      categories: ["Q1", "Q2", "Q3"]
      series:
        - name: "营收"
          values: [10, 20, 15]
'''


class TestP4:
    def test_chart_is_native_ooxml(self, tmp_path):
        """A chart element must export as a real OOXML chart part (selectable
        in PowerPoint), not a flattened image."""
        pytest.importorskip('pptx')
        import zipfile
        from lib.slides.export_pptx import export_pptx
        deck = parse_deck(_write_deck(tmp_path, [_CHART_PAGE]))
        out = str(tmp_path / 'out.pptx')
        export_pptx(deck, out)
        with zipfile.ZipFile(out) as z:
            charts = [n for n in z.namelist()
                      if n.startswith('ppt/charts/chart')]
            assert charts, 'no native chart part written'
            xml = z.read(charts[0]).decode()
        assert '<c:barChart>' in xml
        assert '营收' in xml and 'Q2' in xml

    def test_font_embedding_structure(self, tmp_path):
        """Embedded fonts: fntdata parts + rels + embeddedFontLst placed after
        notesSz (CT_Presentation order), at most one regular and one bold
        slot, and the bytes are glyf-outline TTFs (PowerPoint rejects CFF)."""
        pytest.importorskip('pptx')
        from lib.design_sys import fonts as _fonts
        if not _fonts.ensure_font('misans', 400):
            pytest.skip('misans not staged locally')
        import re
        import zipfile
        from lib.slides.export_pptx import export_pptx
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        out = str(tmp_path / 'out.pptx')
        summary = export_pptx(deck, out)
        assert summary['embeddedFonts'] >= 1
        with zipfile.ZipFile(out) as z:
            assert z.testzip() is None
            parts = [n for n in z.namelist() if n.startswith('ppt/fonts/')]
            assert parts
            pres = z.read('ppt/presentation.xml').decode()
            rels = z.read('ppt/_rels/presentation.xml.rels').decode()
            ct = z.read('[Content_Types].xml').decode()
            blobs = {n: z.read(n) for n in parts}
        assert 'embedTrueTypeFonts="1"' in pres
        assert 'Extension="fntdata"' in ct
        assert '/relationships/font' in rels
        assert pres.index('<p:embeddedFontLst>') > pres.index('<p:notesSz')
        lst = re.search(r'<p:embeddedFontLst>.*?</p:embeddedFontLst>',
                        pres, re.DOTALL).group(0)
        assert lst.count('<p:regular ') <= 1 and lst.count('<p:bold ') <= 1
        # python-pptx must still open the re-zipped package
        from pptx import Presentation
        assert len(Presentation(out).slides) == 1
        # embedded bytes are TrueType-outline sfnt
        from fontTools.ttLib import TTFont
        import io
        for n, blob in blobs.items():
            f = TTFont(io.BytesIO(blob))
            assert 'glyf' in f, f'{n} is not a glyf-outline font'

    def test_import_round_trip_table_text(self, tmp_path):
        """Export → import: table cell text must survive the loop (the DJI
        golden deck page-17 regression — empty cells on first pass)."""
        pytest.importorskip('pptx')
        import glob
        import yaml
        from lib.slides.export_pptx import export_pptx
        from lib.slides.import_pptx import import_pptx
        deck = parse_deck(_write_deck(tmp_path, [_TABLE_PAGE]))
        out = str(tmp_path / 'out.pptx')
        export_pptx(deck, out)
        import_pptx(out, str(tmp_path / 'reimport'))
        found = []
        for p in glob.glob(str(tmp_path / 'reimport' / 'pages' / '*.page')):
            d = yaml.safe_load(open(p, encoding='utf-8'))
            for el in d.get('elements') or []:
                if el.get('elementType') == 'table':
                    for row in el.get('rows') or []:
                        found.append([
                            (c.get('text') if isinstance(c, dict) else c)
                            for c in row])
        assert ['指标', '2025'] in found
        assert ['利润', '15.8'] in found

    def test_edit_slides_tool_schema(self):
        from lib.tools.produce import (EDIT_SLIDES_TOOL,
                                       EDIT_SLIDES_TOOL_NAME,
                                       PRODUCE_TOOL_NAMES)
        assert EDIT_SLIDES_TOOL_NAME in PRODUCE_TOOL_NAMES
        fn = EDIT_SLIDES_TOOL['function']
        assert fn['name'] == 'edit_slides'
        assert set(fn['parameters']['required']) >= {'task_id', 'page',
                                                     'instruction'}
